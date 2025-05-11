# -*- coding: utf-8 -*-
from docx import Document
from docx.shared import Inches, Pt, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import nsdecls
from docx.oxml import parse_xml
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch, cm
from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib.enums import TA_CENTER, TA_LEFT
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak, Image as RLImage, BaseDocTemplate, PageTemplate, Frame
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
import os
import json
import sys
from datetime import datetime
import shutil
import glob
from PIL import Image
import traceback
import logging # Logging modülünü import et
import qrcode
from pptx import Presentation
from pptx.util import Inches as PPTXInches, Pt as PPTXPt # pptx için Inches ve Pt import et (docx ile çakışmaması için alias)

# Logging yapılandırması
logging.basicConfig(level=logging.INFO, stream=sys.stdout,
                    format='%(asctime)s - %(levelname)s - %(message)s')

class DocumentGenerator:
    """
    Arsa analizi verilerinden Word, PDF ve PowerPoint belgeleri oluşturan sınıf.
    """
    def __init__(self, arsa_data, analiz_ozeti, file_id, output_dir, profile_info=None, settings=None):
        """
        DocumentGenerator sınıfını başlatır.

        Args:
            arsa_data (dict): Arsa ile ilgili detaylı veriler.
            analiz_ozeti (dict): Analiz sonuçlarının özeti.
            file_id (str): Oluşturulacak dosyalar için benzersiz kimlik.
            output_dir (str): Oluşturulan dosyaların kaydedileceği ana dizin.
            profile_info (dict, optional): Portföy sorumlusu bilgileri. Defaults to None.
            settings (dict, optional): Belge oluşturma ayarları (tema, renk şeması, bölümler). Defaults to None.
        """
        logging.info("DocumentGenerator __init__ metodu başladı.")
        self.arsa_data = arsa_data or {} # None gelirse boş dict yap
        self.analiz_ozeti = analiz_ozeti or {} # None gelirse boş dict yap
        self.file_id = file_id
        self.output_dir = output_dir

        self.analiz_id = self.arsa_data.get('id')
        logging.info(f"Analiz ID: {self.analiz_id}, File ID: {self.file_id}")

        # Sunum klasörünü output_dir içinde file_id ile oluştur
        self.sunum_klasoru = os.path.join(output_dir, file_id)
        logging.info(f"Sunum Klasörü: {self.sunum_klasoru}")

        try:
            os.makedirs(self.sunum_klasoru, exist_ok=True) # exist_ok=True ile klasör varsa hata vermez
            logging.info(f"Sunum klasörü oluşturuldu veya zaten var: {self.sunum_klasoru}")
        except OSError as e:
            logging.error(f"Sunum klasörü oluşturulamadı: {e}")
            # Hata durumunda daha sonra dosya yazma işlemleri başarısız olacaktır.
            # Bu hatayı fırlatmak veya daha sonraki işlemleri kontrol etmek gerekebilir.
            # Şimdilik logluyoruz ve devam ediyoruz, ancak dosya yazma hataları beklenir.


        self.profile_info = profile_info or {}
        # Ayarlardan bağımsız olarak tüm bölümleri varsayılan olarak dahil et
        self.settings = settings or {
            'theme': 'classic', # Gelecekte tema desteği için
            'color_scheme': 'blue',
            # 'sections' listesi artık _should_include_section tarafından kullanılmayacak
            'sections': ['profile', 'property', 'infrastructure', 'swot', 'photos', 'summary']
        }
        self.colors = self._get_color_scheme()

        # Logo yolu artık daha dinamik olabilir veya yapılandırmadan alınabilir.
        # Şimdilik sabit bırakalım ama app.py'nin bir üst dizinindeki static'e işaret etmeli.
        base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__))) # modules'ten app.py'nin olduğu yere
        self.logo_path = os.path.join(base_dir, 'static', 'logo.png')
        logging.info(f"Logo yolu kontrol ediliyor: {self.logo_path}")


    def _get_color_scheme(self):
        """Seçilen renk şemasına göre renkleri döndürür."""
        schemes = {
            'blue': {
                'primary': RGBColor(26, 35, 126),  # Koyu Mavi
                'secondary': RGBColor(63, 81, 181), # Orta Mavi
                'accent': RGBColor(33, 150, 243),   # Açık Mavi
                'background_light': 'E3F2FD',       # En Açık Mavi (PDF için hex)
                'background_medium': 'E8EAF6'       # Biraz daha koyu mavi (Word için hex)
            },
            'green': {
                'primary': RGBColor(27, 94, 32),    # Koyu Yeşil
                'secondary': RGBColor(56, 142, 60),  # Orta Yeşil
                'accent': RGBColor(76, 175, 80),     # Açık Yeşil
                'background_light': 'E8F5E9',
                'background_medium': 'DCEDC8'
            },
            'purple': {
                'primary': RGBColor(74, 20, 140),    # Koyu Mor
                'secondary': RGBColor(123, 31, 162), # Orta Mor
                'accent': RGBColor(156, 39, 176),    # Açık Mor
                'background_light': 'F3E5F5',
                'background_medium': 'E1BEE7'
            }
        }
        # Belirtilen renk şeması yoksa varsayılan olarak 'blue' kullanılır.
        return schemes.get(self.settings.get('color_scheme', 'blue'), schemes['blue'])

    def _should_include_section(self, section):
        """Belirli bir bölümün rapora dahil edilip edilmeyeceğini kontrol eder."""
        # Bu metod artık kullanılmayacak, tüm bölümler varsayılan olarak dahil ediliyor.
        # Gelecekte bölüm seçimi eklendiğinde bu metod tekrar aktif edilebilir.
        # return section in self.settings.get('sections', []) # 'sections' yoksa boş liste döndür
        return True # Tüm bölümleri her zaman dahil et


    def _format_currency(self, value):
        """Para birimini formatlayan yardımcı metod."""
        try:
            float_value = float(value)
            # Türkçe format: binlik ayıracı '.', ondalık ayıracı ','
            return f"{float_value:,.2f} TL".replace(',', 'X').replace('.', ',').replace('X', '.')
        except (TypeError, ValueError):
            return "0,00 TL"

    def _format_area(self, value):
        """Alan birimini formatlayan yardımcı metod."""
        try:
            float_value = float(value)
            # Türkçe format: binlik ayıracı '.', ondalık ayıracı ','
            return f"{float_value:,.2f} m²".replace(',', 'X').replace('.', ',').replace('X', '.')
        except (TypeError, ValueError):
            return "0,00 m²"

    def _get_arsa_bilgileri(self):
        """Arsa bilgilerini tablo formatında hazırlar."""
        # Eksik veya None değerler için güvenli get kullanımı
        il = self.arsa_data.get('il', '')
        ilce = self.arsa_data.get('ilce', '')
        mahalle = self.arsa_data.get('mahalle', '')
        ada = self.arsa_data.get('ada', '')
        parsel = self.arsa_data.get('parsel', '')
        imar_durumu = self.arsa_data.get('imar_durumu', '')
        taks = self.arsa_data.get('taks', '')
        kaks = self.arsa_data.get('kaks', '')

        try:
            fiyat = float(self.arsa_data.get('fiyat', 0))
        except (TypeError, ValueError):
            fiyat = 0
            logging.warning("Arsa fiyatı geçerli sayı formatında değil.")

        try:
            metrekare = float(self.arsa_data.get('metrekare', 1))
        except (TypeError, ValueError):
            metrekare = 1
            logging.warning("Arsa metrekare geçerli sayı formatında değil.")

        metrekare_fiyati = fiyat / metrekare if metrekare > 0 else 0

        return [
            ['İl/İlçe', f"{il}/{ilce}"],
            ['Mahalle', mahalle],
            ['Ada/Parsel', f"{ada}/{parsel}"],
            ['Alan', self._format_area(metrekare)],
            ['İmar Durumu', imar_durumu],
            ['TAKS/KAKS', f"{taks}/{kaks}"],
            ['Toplam Fiyat', self._format_currency(fiyat)],
            ['m² Fiyatı', self._format_currency(metrekare_fiyati)],
        ]

    def _get_altyapi_durumu(self):
        """Altyapı durumunu tablo formatında hazırlar."""
        # app.py'den gelen 'altyapi[]' anahtarını doğru şekilde al
        altyapi_list = self.arsa_data.get('altyapi[]', [])

        # Gelen veri string ise listeye çevir (tek elemanlı)
        if isinstance(altyapi_list, str):
            altyapi_list = [altyapi_list]
        # Gelen veri None ise boş liste yap
        elif altyapi_list is None:
            altyapi_list = []
        # Gelen veri liste değilse (beklenmeyen durum), boş liste yap
        elif not isinstance(altyapi_list, list):
            logging.warning(f"Beklenmeyen altyapi[] veri tipi: {type(altyapi_list)}")
            altyapi_list = []


        return [
            ['Yol', '✓' if 'yol' in altyapi_list else '✗'],
            ['Elektrik', '✓' if 'elektrik' in altyapi_list else '✗'],
            ['Su', '✓' if 'su' in altyapi_list else '✗'],
            ['Doğalgaz', '✓' if 'dogalgaz' in altyapi_list else '✗'],
            ['Kanalizasyon', '✓' if 'kanalizasyon' in altyapi_list else '✗']
        ]

    def _get_swot_analizi(self):
        """SWOT analizini dict formatında hazırlar."""
        swot_data = {
            'Güçlü Yönler': self.arsa_data.get('strengths', []),
            'Zayıf Yönler': self.arsa_data.get('weaknesses', []),
            'Fırsatlar': self.arsa_data.get('opportunities', []),
            'Tehditler': self.arsa_data.get('threats', [])
        }

        # Gelen verinin liste olduğundan emin ol
        for key, value in swot_data.items():
            if isinstance(value, str):
                try:
                    # JSON string ise parse et
                    parsed_value = json.loads(value)
                    if isinstance(parsed_value, list):
                        swot_data[key] = parsed_value
                    else:
                        # JSON ama liste değilse, tek elemanlı listeye çevir
                        swot_data[key] = [str(parsed_value)]
                        logging.warning(f"SWOT verisi JSON ancak liste değil: {key}")
                except json.JSONDecodeError:
                    # JSON değilse ve boş değilse, tek elemanlı listeye çevir
                    swot_data[key] = [value] if value else []
                    logging.warning(f"SWOT verisi JSON formatında değil: {key}")
            elif not isinstance(value, list):
                 # Liste veya string değilse (örn. None), boş listeye çevir
                 swot_data[key] = []
                 logging.warning(f"SWOT verisi beklenmeyen formatta: {key}, Tipi: {type(value)}")

        return swot_data

    def _set_cell_background(self, cell, color_hex):
        """Word tablo hücresinin arka plan rengini ayarlar."""
        try:
            shading_elm = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{color_hex}"/>')
            cell._tc.get_or_add_tcPr().append(shading_elm)
        except Exception as e:
            logging.error(f"Hücre arkaplanı ayarlanamadı: {e}")

    def _get_uploaded_images(self):
        """Yüklenen resimlerin yollarını bulur ve sunum klasörüne kopyalar."""
        logging.info("_get_uploaded_images metodu başladı.")
        if not self.analiz_id:
            logging.warning("Analiz ID bulunamadı. Resimler aranamıyor.")
            return []

        analiz_id_str = str(self.analiz_id)
        # app.py'nin bulunduğu dizine göre static/uploads klasörünü bul
        base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        # Düzeltilen yol: static/uploads/{analiz_id}
        uploads_dir = os.path.join(base_dir, 'static', 'uploads', analiz_id_str)
        logging.info(f"Orijinal resimlerin aranacağı klasör: {uploads_dir}")

        image_extensions = ['*.jpg', '*.jpeg', '*.png', '*.gif']
        original_images = []
        if os.path.exists(uploads_dir):
            for ext in image_extensions:
                search_pattern = os.path.join(uploads_dir, ext)
                found = glob.glob(search_pattern)
                if found:
                     original_images.extend(found)
        else:
            logging.warning(f"Yükleme klasörü bulunamadı: {uploads_dir}")

        original_images = sorted(list(set(original_images))) # Tekrar edenleri kaldır ve sırala
        logging.info(f"Bulunan toplam orijinal resim sayısı: {len(original_images)}")

        copied_images = []
        logging.info(f"Kopyalama hedef klasörü: {self.sunum_klasoru}")
        for img_path in original_images:
            dest_path = os.path.join(self.sunum_klasoru, os.path.basename(img_path))
            try:
                shutil.copy(img_path, dest_path)
                copied_images.append(dest_path)
                logging.info(f"Resim kopyalandı: {img_path} -> {dest_path}")
            except Exception as e:
                logging.error(f"Resim kopyalanamadı! Kaynak: {img_path}, Hedef: {dest_path}, Hata: {e}")
                # Hata durumunda bu resim listeye eklenmez.

        logging.info(f"Kopyalanan resim sayısı: {len(copied_images)}")
        return copied_images

    def _get_profile_table_data(self):
        """Profil bilgilerini tablo formatında hazırlar."""
        p = self.profile_info or {}
        created_at = p.get('created_at')
        # created_at bir datetime nesnesi olabilir, string olabilir veya None olabilir.
        if isinstance(created_at, datetime):
             date_str = created_at.strftime('%d.%m.%Y %H:%M')
        elif isinstance(created_at, str) and created_at:
             # String ise parse etmeyi deneyebiliriz veya direkt kullanabiliriz
             try:
                 dt_obj = datetime.fromisoformat(created_at) # ISO formatı bekleniyor
                 date_str = dt_obj.strftime('%d.%m.%Y %H:%M')
             except ValueError:
                 date_str = created_at # Parse edilemezse string olarak kullan
                 logging.warning(f"Analiz tarihi string formatında parse edilemedi: {created_at}")
        else:
            date_str = '-'


        return [
            ['Ad Soyad', f"{p.get('ad','') or ''} {p.get('soyad','') or ''}"],
            ['Ünvan', p.get('unvan', '-')],
            ['Firma', p.get('firma', '-')],
            ['E-posta', p.get('email', '-')],
            ['Telefon', p.get('telefon', '-')],
            ['Adres', p.get('adres', '-')],
            ['Analiz Tarihi', date_str]
        ]

    def _get_profile_photo_path(self):
        """Profil fotoğrafı varsa tam yolunu döndürür."""
        pf = self.profile_info.get('profil_foto')
        if pf:
            # pf ör: 'profiles/1/xxx.jpg'
            # app.py'nin bulunduğu dizine göre static/uploads klasörünü bul
            base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
            # Düzeltilen yol: static/uploads/{profil_foto yolu}
            # pf zaten profiles/1/xxx.jpg gibi bir yol içeriyor, bunu static/uploads altına eklemeliyiz
            # Ancak profile_info'daki 'profil_foto' alanı doğrudan 'uploads' klasörünün altındaki yolu tutuyorsa
            # (örneğin 'profiles/1/xxx.jpg'), o zaman static'i eklememiz gerekir.
            # Logdaki hata '...uploads\profiles/1/IMG-20250425-WA0001.jpg' gösteriyor, bu da 'uploads' altında 'profiles'
            # klasörünün olduğunu düşündürüyor. Bu durumda static'i ekleyelim.
            uploads_dir = os.path.join(base_dir, 'static', 'uploads')
            full_path = os.path.join(uploads_dir, pf.replace('/', os.sep)) # Yol ayırıcıyı düzelt
            if os.path.exists(full_path):
                return full_path
            else:
                logging.warning(f"Profil fotoğrafı dosyası bulunamadı: {full_path}")
        return None

    def _add_footer_word(self, doc, text):
        """Tüm Word sayfalarına footer ekler."""
        try:
            for section in doc.sections:
                # Footer'ı temizle veya ilk paragrafı al
                if not section.footer.paragraphs:
                    footer = section.footer.add_paragraph()
                else:
                    footer = section.footer.paragraphs[0]

                footer.text = text
                footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
                # Mevcut run'ları temizleyip yeniden oluşturmak daha güvenli olabilir
                # Ancak şimdilik text'i set etmek yeterli.
                if footer.runs:
                    run = footer.runs[0]
                else:
                    run = footer.add_run()
                    run.text = text # Tekrar text atama

                run.font.size = Pt(9)
                run.font.name = 'Arial' # Footer için sabit font
                run.font.color.rgb = RGBColor(120, 120, 120)
        except Exception as e:
            logging.error(f"Word footer eklenirken hata oluştu: {e}")


    def _add_footer_pdf(self, canvas, doc):
        """PDF için footer fonksiyonu."""
        try:
            canvas.saveState()
            canvas.setFont('Helvetica', 9) # PDF footer için sabit font
            canvas.setFillColor(colors.HexColor('#888888'))
            canvas.drawString(2*cm, 1*cm, "Gayrimenkul Analiz Sistemi - www.invecoproje.com")
            canvas.drawRightString(A4[0]-2*cm, 1*cm, f"Sayfa {doc.page}")
            canvas.restoreState()
        except Exception as e:
            logging.error(f"PDF footer eklenirken hata oluştu: {e}")


    def _swot_table_data(self):
        """SWOT'u tabloya uygun şekilde hazırlar (liste içinde liste)."""
        swot = self._get_swot_analizi()
        # Her maddeyi ayrı bir satır/paragraf yapmak için formatla
        formatted_swot = {
            key: [f"• {item}" for item in items] if items else ["- Yok -"]
            for key, items in swot.items()
        }
        return [
            [formatted_swot.get('Güçlü Yönler', ["- Yok -"]), formatted_swot.get('Zayıf Yönler', ["- Yok -"])],
            [formatted_swot.get('Fırsatlar', ["- Yok -"]), formatted_swot.get('Tehditler', ["- Yok -"])]
        ]

    def _get_insaat_hesaplama(self):
        """İnşaat alanı hesaplama verilerini tablo formatında hazırlar."""
        # ArsaAnalizci kullanılmadığı için, arsa_data'dan değerleri çek
        hesap = self.arsa_data.get('insaat_hesaplama')

        # Eğer insaat_hesaplama yoksa veya eksikse, manuel hesapla
        if not hesap or not all(k in hesap for k in ['taban_alani', 'toplam_insaat_alani', 'teorik_kat_sayisi', 'tam_kat_sayisi']):
            logging.warning("İnşaat hesaplama verisi eksik veya hatalı, manuel hesaplanıyor.")
            try:
                metrekare = float(self.arsa_data.get('metrekare', 0))
                # TAKS/KAKS değerlerini float'a çevirirken hata yönetimi
                try:
                    taks = float(self.arsa_data.get('taks', 0.3)) # Varsayılan değer
                except (TypeError, ValueError):
                    taks = 0.3
                    logging.warning("TAKS değeri geçerli sayı formatında değil, varsayılan 0.3 kullanıldı.")

                try:
                    kaks = float(self.arsa_data.get('kaks', 1.5)) # Varsayılan değer
                except (TypeError, ValueError):
                    kaks = 1.5
                    logging.warning("KAKS değeri geçerli sayı formatında değil, varsayılan 1.5 kullanıldı.")

                taban_alani = metrekare * taks
                toplam_insaat_alani = metrekare * kaks
                teorik_kat_sayisi = kaks / taks if taks and taks > 0 else 0 # Sıfıra bölme kontrolü
                tam_kat_sayisi = int(teorik_kat_sayisi) # Tam sayıya çevir

                hesap = {
                    'taban_alani': taban_alani,
                    'toplam_insaat_alani': toplam_insaat_alani,
                    'teorik_kat_sayisi': teorik_kat_sayisi,
                    'tam_kat_sayisi': tam_kat_sayisi
                }
            except Exception as e:
                logging.error(f"Manuel inşaat hesaplama hatası: {e}")
                hesap = { # Hata durumunda varsayılan değerler
                    'taban_alani': 0,
                    'toplam_insaat_alani': 0,
                    'teorik_kat_sayisi': 0,
                    'tam_kat_sayisi': 0
                }
        else:
             logging.info("İnşaat hesaplama verisi arsa_data'dan alındı.")


        return [
            ['Taban Alanı (m²)', f"{hesap.get('taban_alani', 0):.2f}"],
            ['Toplam İnşaat Alanı (m²)', f"{hesap.get('toplam_insaat_alani', 0):.2f}"],
            ['Teorik Kat Sayısı', f"{hesap.get('teorik_kat_sayisi', 0):.2f}"],
            ['Tam Kat Sayısı', f"{hesap.get('tam_kat_sayisi', 0)}"]
        ]

    def _create_qr_code(self):
        """Analiz detayları için QR kod oluşturur ve kaydeder."""
        try:
            # QR kodun içeriği - analiz detay sayfasının URL'i
            # URL'i yapılandırmadan veya dinamik olarak almak daha iyi olabilir.
            # Şimdilik sabit IP kullanılıyor, bu prod ortamında değişmeli.
            qr_data = f"http://192.168.2.14:5000/analiz/{self.arsa_data.get('id')}"
            if not self.arsa_data.get('id'):
                 logging.warning("Analiz ID olmadığı için QR kod içeriği eksik olacak.")
                 qr_data = "Analiz ID bulunamadı."

            # QR kod ayarları
            qr = qrcode.QRCode(
                version=1,
                error_correction=qrcode.constants.ERROR_CORRECT_L,
                box_size=10,
                border=4,
            )
            qr.add_data(qr_data)
            qr.make(fit=True)

            # QR kod görüntüsü oluştur
            qr_image = qr.make_image(fill_color="black", back_color="white")

            # QR kodu kaydet
            qr_path = os.path.join(self.sunum_klasoru, 'qr_code.png')
            qr_image.save(qr_path)
            logging.info(f"QR kod oluşturuldu ve kaydedildi: {qr_path}")

            return qr_path
        except Exception as e:
            logging.error(f"QR kod oluşturma hatası: {e}")
            return None

    def create_word(self):
        """Word belgesi oluşturur."""
        logging.info("Word belgesi oluşturma başladı.")
        try:
            doc = Document()

            # Sayfa yapılandırması (A4 Yatay)
            # İlk section zaten var, diğerlerini eklemeye gerek yok
            section = doc.sections[0]
            section.orientation = 1 # WD_ORIENT.LANDSCAPE (enum yerine direkt değer)
            section.page_width = Cm(29.7)
            section.page_height = Cm(21.0)
            section.left_margin = Cm(1.27) # 0.5 inch
            section.right_margin = Cm(1.27)
            section.top_margin = Cm(1.27)
            section.bottom_margin = Cm(1.27)

            # Kapak Sayfası
            self._add_word_cover_page(doc)
            doc.add_page_break()

            # PROFİL SAYFASI
            self._add_word_profile_page(doc)
            doc.add_page_break()

            # İçindekiler Sayfası (manuel olarak eklendi, otomatik TOC daha karmaşık)
            self._add_word_toc_page(doc)
            doc.add_page_break()

            # Arsa Bilgileri Sayfası
            self._add_word_property_page(doc)
            doc.add_page_break()

            # Altyapı Durumu Sayfası
            self._add_word_infrastructure_page(doc)
            doc.add_page_break()

            # İnşaat Alanı Hesaplama Sayfası
            self._add_word_insaat_hesaplama_page(doc)
            doc.add_page_break()

            # SWOT Analizi Sayfası
            self._add_word_swot_page(doc)
            doc.add_page_break()

            # Analiz Özeti Sayfası
            self._add_word_summary_page(doc)
            doc.add_page_break()

            # Arsa Fotoğrafları Bölümü
            self._add_word_photos_page(doc)
            doc.add_page_break() # Son bölümden sonra sayfa sonu ekle

            # QR KOD Sayfası
            qr_path = self._create_qr_code()
            if qr_path and os.path.exists(qr_path):
                 self._add_word_qr_code_page(doc, qr_path)
                 doc.add_page_break() # Son bölümden sonra sayfa sonu ekle


            # Footer ekle
            self._add_footer_word(doc, "Gayrimenkul Analiz Sistemi - invecoproje.com | " + datetime.now().strftime('%d.%m.%Y'))

            # Generate the filename
            filename = os.path.join(self.sunum_klasoru, f'analiz_{self.file_id}.docx')

            # Save the document
            doc.save(filename)
            logging.info(f"Word dosyası başarıyla kaydedildi: {filename}")

            # Verify the file was created
            if os.path.exists(filename):
                return filename  # Return the full path if file exists
            else:
                logging.error(f"Word dosyası oluşturulamadı: {filename}")
                return None

        except Exception as e:
            logging.error(f"Word belgesi oluşturulurken genel hata: {e}")
            traceback.print_exc(file=sys.stdout)
            sys.stdout.flush()
            return None

    def _add_word_cover_page(self, doc):
        """Word belgesine kapak sayfası ekler."""
        # Logo ekle (varsa)
        if os.path.exists(self.logo_path):
            try:
                p_logo = doc.add_paragraph()
                p_logo.alignment = WD_ALIGN_PARAGRAPH.CENTER
                run_logo = p_logo.add_run()
                # Logo boyutunu ayarla
                run_logo.add_picture(self.logo_path, width=Inches(2.5))
                p_logo.space_after = Pt(20) # Logo sonrası boşluk
            except Exception as e:
                logging.error(f"Word kapak sayfasına logo eklenemedi: {e}")

        # Başlık
        title = doc.add_heading('', 0)
        title_run = title.add_run('GAYRİMENKUL ANALİZ RAPORU')
        title_run.font.name = 'Arial'
        title_run.font.size = Pt(40)
        title_run.font.color.rgb = self.colors['primary']
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        title.space_after = Pt(10) # Başlık sonrası boşluk

        # Alt başlık (İl/İlçe)
        subtitle = doc.add_paragraph()
        subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
        subtitle_run = subtitle.add_run(f"{self.arsa_data.get('il', '')}, {self.arsa_data.get('ilce', '')}")
        subtitle_run.font.name = 'Arial'
        subtitle_run.font.size = Pt(28)
        subtitle_run.font.color.rgb = self.colors['secondary']
        subtitle.space_before = Pt(12)
        subtitle.space_after = Pt(40) # Alt başlık sonrası boşluk

        # Tarih
        date_paragraph = doc.add_paragraph()
        date_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        date_run = date_paragraph.add_run(datetime.now().strftime('%d.%m.%Y'))
        date_run.font.name = 'Arial'
        date_run.font.size = Pt(14)
        date_run.font.color.rgb = RGBColor(96, 125, 139) # Gri Mavi

    def _add_word_profile_page(self, doc):
        """Word belgesine profil bilgileri sayfası ekler."""
        heading = doc.add_heading('Portföy Sorumlusu', 1)
        heading.runs[0].font.name = 'Arial'
        heading.runs[0].font.size = Pt(22)
        heading.runs[0].font.color.rgb = self.colors['accent']
        heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
        heading.space_after = Pt(20) # Başlık sonrası boşluk

        # Profil fotoğrafı (varsa)
        profile_photo_path = self._get_profile_photo_path()
        if profile_photo_path:
            try:
                p_photo = doc.add_paragraph()
                p_photo.alignment = WD_ALIGN_PARAGRAPH.CENTER
                run_photo = p_photo.add_run()
                run_photo.add_picture(profile_photo_path, width=Inches(1.8)) # Fotoğraf boyutunu ayarla
                p_photo.space_after = Pt(15) # Fotoğraf sonrası boşluk
            except Exception as e:
                logging.error(f"Word profil sayfasına fotoğraf eklenemedi: {e}")


        # Profil bilgileri tablosu
        profile_data = self._get_profile_table_data()
        table = doc.add_table(rows=len(profile_data), cols=2)
        table.style = 'Table Grid'
        table.autofit = False
        table.allow_autofit = False

        # Sütun genişliklerini ayarla (sayfa genişliğine göre)
        section = doc.sections[0]
        available_width = section.page_width - section.left_margin - section.right_margin
        col1_width = available_width * 0.35
        col2_width = available_width * 0.65

        for i, (label, value) in enumerate(profile_data):
            cell1 = table.cell(i, 0)
            cell2 = table.cell(i, 1)

            # Sütun genişliklerini set et
            cell1.width = col1_width
            cell2.width = col2_width


            p1 = cell1.paragraphs[0]
            p1.clear() # Mevcut paragrafı temizle
            run1 = p1.add_run(label)
            run1.font.bold = True
            run1.font.name = 'Arial'
            run1.font.size = Pt(11) # Font boyutu
            p1.alignment = WD_ALIGN_PARAGRAPH.LEFT
            self._set_cell_background(cell1, self.colors['background_light']) # Açık mavi/tema rengi

            p2 = cell2.paragraphs[0]
            p2.clear() # Mevcut paragrafı temizle
            run2 = p2.add_run(value)
            run2.font.name = 'Arial'
            run2.font.size = Pt(11) # Font boyutu
            p2.alignment = WD_ALIGN_PARAGRAPH.LEFT
            self._set_cell_background(cell2, "FFFFFF") # Beyaz

            cell1.vertical_alignment = 1 # Dikey ortala
            cell2.vertical_alignment = 1 # Dikey ortala

        # Tablo sonrası boşluk
        doc.add_paragraph().space_after = Pt(10)

    def _add_word_toc_page(self, doc):
        """Word belgesine içindekiler sayfası ekler."""
        toc_heading = doc.add_heading('İçindekiler', level=1)
        toc_heading.runs[0].font.name = 'Arial'
        toc_heading.runs[0].font.size = Pt(24)
        toc_heading.runs[0].font.color.rgb = self.colors['primary']
        toc_heading.space_after = Pt(20) # Başlık sonrası boşluk

        # İçindekiler listesi (manuel) - Tüm bölümleri ekleyelim
        sections_list = [
            'Portföy Sorumlusu',
            'Arsa Bilgileri',
            'Altyapı Durumu',
            'İnşaat Alanı Hesaplaması',
            'SWOT Analizi',
            'Analiz Özeti',
            'Arsa Fotoğrafları',
            'Analiz Detayları QR Kod' # QR kod sayfası da eklendi
        ]

        for i, section_name in enumerate(sections_list, 1):
            p = doc.add_paragraph(style='List Number') # Numaralı liste stili
            p.paragraph_format.left_indent = Inches(0.5)
            run = p.add_run(f"{section_name}")
            run.font.name = 'Arial'
            run.font.size = Pt(14)
            p.space_after = Pt(5) # Madde sonrası boşluk


    def _add_word_property_page(self, doc):
        """Word belgesine arsa bilgileri sayfası ekler."""
        heading = doc.add_heading('Arsa Bilgileri', 1)
        heading.runs[0].font.name = 'Arial'
        heading.runs[0].font.size = Pt(24)
        heading.runs[0].font.color.rgb = self.colors['primary']
        heading.space_after = Pt(20) # Başlık sonrası boşluk

        # Ana bilgi tablosu - 2 sütunlu daha sade tasarım
        data = self._get_arsa_bilgileri()

        table = doc.add_table(rows=len(data), cols=2)
        table.style = 'Table Grid'
        table.autofit = False
        table.allow_autofit = False

        # Sütun genişliklerini ayarla (sayfa genişliğine göre)
        section = doc.sections[0]
        available_width = section.page_width - section.left_margin - section.right_margin
        col1_width = available_width * 0.3
        col2_width = available_width * 0.7

        for i, row_data in enumerate(data):
            cell1 = table.cell(i, 0)
            cell2 = table.cell(i, 1)

            # Sütun genişliklerini set et
            cell1.width = col1_width
            cell2.width = col2_width

            # Özellik hücresi (Sola Yaslı, Kalın)
            p1 = cell1.paragraphs[0]
            p1.clear()
            run1 = p1.add_run(row_data[0])
            run1.font.bold = True
            run1.font.name = 'Arial'
            run1.font.size = Pt(11)
            p1.alignment = WD_ALIGN_PARAGRAPH.LEFT
            self._set_cell_background(cell1, self.colors['background_medium']) # Orta mavi/tema rengi

            # Değer hücresi (Sola Yaslı)
            p2 = cell2.paragraphs[0]
            p2.clear()
            run2 = p2.add_run(str(row_data[1])) # Değeri stringe çevir
            run2.font.name = 'Arial'
            run2.font.size = Pt(11)
            p2.alignment = WD_ALIGN_PARAGRAPH.LEFT
            self._set_cell_background(cell2, "FFFFFF") # Beyaz

            # Hücre dikey hizalama (Ortala)
            cell1.vertical_alignment = 1
            cell2.vertical_alignment = 1

        # Tablo sonrası boşluk
        doc.add_paragraph().space_after = Pt(10)

    def _add_word_infrastructure_page(self, doc):
        """Word belgesine altyapı durumu sayfası ekler."""
        heading = doc.add_heading('Altyapı Durumu', 1)
        heading.runs[0].font.name = 'Arial'
        heading.runs[0].font.size = Pt(24)
        heading.runs[0].font.color.rgb = self.colors['primary']
        heading.space_after = Pt(20) # Başlık sonrası boşluk


        altyapi_data = self._get_altyapi_durumu()
        table = doc.add_table(rows=len(altyapi_data), cols=2)
        table.style = 'Table Grid'
        table.autofit = False
        table.allow_autofit = False

        # Sütun genişliklerini ayarla
        section = doc.sections[0]
        available_width = section.page_width - section.left_margin - section.right_margin
        col1_width = available_width * 0.4
        col2_width = available_width * 0.6


        for i, (altyapi, durum) in enumerate(altyapi_data):
            cell1 = table.cell(i, 0)
            cell2 = table.cell(i, 1)

            # Sütun genişliklerini set et
            cell1.width = col1_width
            cell2.width = col2_width

            # Özellik hücresi
            p1 = cell1.paragraphs[0]
            p1.clear()
            run1 = p1.add_run(altyapi)
            run1.font.bold = True
            run1.font.name = 'Arial'
            run1.font.size = Pt(11)
            p1.alignment = WD_ALIGN_PARAGRAPH.LEFT
            self._set_cell_background(cell1, self.colors['background_medium'])

            # Durum hücresi (İkon ve Renk)
            p2 = cell2.paragraphs[0]
            p2.clear()
            status_text = '✓ Var' if durum == '✓' else '✗ Yok'
            status_run = p2.add_run(status_text)
            status_run.font.name = 'Arial'
            status_run.font.size = Pt(11)
            status_run.font.bold = True
            # Duruma göre renk
            status_run.font.color.rgb = RGBColor(76, 175, 80) if durum == '✓' else RGBColor(244, 67, 54) # Yeşil / Kırmızı
            p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
            self._set_cell_background(cell2, "FFFFFF")

            cell1.vertical_alignment = 1
            cell2.vertical_alignment = 1

        # Tablo sonrası boşluk
        doc.add_paragraph().space_after = Pt(10)


    def _add_word_insaat_hesaplama_page(self, doc):
        """Word belgesine inşaat alanı hesaplama sayfası ekler."""
        heading = doc.add_heading('İnşaat Alanı Hesaplaması', 1)
        heading.runs[0].font.name = 'Arial'
        heading.runs[0].font.size = Pt(20)
        heading.runs[0].font.color.rgb = self.colors['accent']
        heading.space_after = Pt(20) # Başlık sonrası boşluk

        data = self._get_insaat_hesaplama()
        table = doc.add_table(rows=len(data), cols=2)
        table.style = 'Table Grid'
        table.autofit = False
        table.allow_autofit = False

        # Sütun genişliklerini ayarla
        section = doc.sections[0]
        available_width = section.page_width - section.left_margin - section.right_margin
        col1_width = available_width * 0.4
        col2_width = available_width * 0.6


        for i, (label, value) in enumerate(data):
            cell1 = table.cell(i, 0)
            cell2 = table.cell(i, 1)

            # Sütun genişliklerini set et
            cell1.width = col1_width
            cell2.width = col2_width

            p1 = cell1.paragraphs[0]
            p1.clear()
            run1 = p1.add_run(label)
            run1.font.bold = True
            run1.font.name = 'Arial'
            run1.font.size = Pt(11)
            p1.alignment = WD_ALIGN_PARAGRAPH.LEFT
            self._set_cell_background(cell1, self.colors['background_light'])

            p2 = cell2.paragraphs[0]
            p2.clear()
            run2 = p2.add_run(str(value)) # Değeri stringe çevir
            run2.font.name = 'Arial'
            run2.font.size = Pt(11)
            p2.alignment = WD_ALIGN_PARAGRAPH.LEFT
            self._set_cell_background(cell2, "FFFFFF")

            cell1.vertical_alignment = 1
            cell2.vertical_alignment = 1

        # Tablo sonrası boşluk
        doc.add_paragraph().space_after = Pt(10)


    def _add_word_swot_page(self, doc):
        """Word belgesine SWOT analizi sayfası ekler."""
        swot_heading = doc.add_heading('SWOT Analizi', 1)
        swot_heading.runs[0].font.name = 'Arial'
        swot_heading.runs[0].font.size = Pt(24)
        swot_heading.runs[0].font.color.rgb = self.colors['primary']
        swot_heading.space_after = Pt(20) # Başlık sonrası boşluk


        swot_titles = ['Güçlü Yönler', 'Zayıf Yönler', 'Fırsatlar', 'Tehditler']
        # Renkler Hex formatında olmalı
        swot_colors_hex = ['A5D6A7', 'EF9A9A', '90CAF9', 'FFD54F'] # Yeşil, Kırmızı, Mavi, Turuncu tonları
        swot_data = self._swot_table_data() # Formatlanmış madde listeleri

        table = doc.add_table(rows=2, cols=2)
        table.style = 'Table Grid'
        table.allow_autofit = False

        # Sütun genişliklerini ayarla
        section = doc.sections[0]
        available_width = section.page_width - section.left_margin - section.right_margin
        col_width = available_width / 2

        for i in range(2):
            for j in range(2):
                idx = i*2 + j
                cell = table.cell(i, j)
                cell.width = col_width # Sütun genişliğini set et

                # Başlık ve maddeler için tek bir paragraf kullanmak yerine
                # başlık için ayrı, maddeler için ayrı paragraflar ekleyelim.
                p_title = cell.paragraphs[0]
                p_title.clear()
                run_title = p_title.add_run(swot_titles[idx])
                run_title.font.bold = True
                run_title.font.size = Pt(13)
                run_title.font.name = 'Arial'
                # Başlık rengini ayarlamak için run'a renk ataması yapabiliriz,
                # ancak hücre arkaplan rengi zaten tema rengiyle ayarlanıyor.
                # İsteğe bağlı olarak başlık run'ına da renk atayabiliriz.
                # run_title.font.color.rgb = RGBColor.from_string(swot_colors_hex[idx]) # Başlık rengini de ayarla

                self._set_cell_background(cell, swot_colors_hex[idx]) # Hücre arkaplan rengi
                p_title.alignment = WD_ALIGN_PARAGRAPH.CENTER # Başlığı ortala
                p_title.space_after = Pt(6) # Başlık sonrası boşluk

                # Maddeler
                items = swot_data[i][j] # Formatlanmış madde listesi
                if items:
                    for madde_idx, madde_text in enumerate(items):
                        # Her madde için yeni bir paragraf oluştur, ancak ilk madde için mevcut olanı kullanma
                        # Bu yaklaşım, List Bullet stilini doğru uygulamak için daha iyi olabilir.
                        # Eğer cell.add_paragraph() kullanırsak, ilk başlık paragrafı kaybolabilir.
                        # Bu yüzden, maddeleri başlık paragrafından sonra ekleyelim.
                        if madde_idx == 0 and not p_title.text: # Eğer başlık paragrafı boşsa (clear sonrası)
                             p_item = p_title # İlk madde için başlık paragrafını kullan
                        else:
                             p_item = cell.add_paragraph()

                        # Madde işaretli liste stili uygulamak için stil adı kullanılabilir
                        # veya manuel olarak madde işareti eklenebilir.
                        # Python-docx'te stil uygulamak daha doğru olur.
                        # Eğer 'List Bullet' stili belgede tanımlıysa:
                        try:
                            p_item.style = 'List Bullet'
                        except KeyError: # Stil bulunamazsa manuel ekle
                            run_item_bullet = p_item.add_run("• ")
                            run_item_bullet.font.name = 'Arial'
                            run_item_bullet.font.size = Pt(10)

                        run_item = p_item.add_run(madde_text.replace("• ","").strip()) # Madde işaretini tekrar eklememek için
                        run_item.font.name = 'Arial'
                        run_item.font.size = Pt(10)
                        p_item.paragraph_format.left_indent = Inches(0.2) # Girinti
                        p_item.paragraph_format.space_before = Pt(0)
                        p_item.paragraph_format.space_after = Pt(3)

                else:
                    p_item = cell.add_paragraph() # Başlıktan sonra yeni paragraf
                    run_item = p_item.add_run("- Yok -")
                    run_item.font.name = 'Arial'
                    run_item.font.size = Pt(10)
                    run_item.italic = True # Yoksa italik yap
                    p_item.alignment = WD_ALIGN_PARAGRAPH.LEFT


                cell.vertical_alignment = 0 # Dikeyde üste hizala

        # Tablo sonrası boşluk
        doc.add_paragraph().space_after = Pt(10)

    def _add_word_summary_page(self, doc):
        """Word belgesine analiz özeti sayfası ekler."""
        ozet_heading = doc.add_heading('Analiz Özeti', 1)
        ozet_heading.runs[0].font.name = 'Arial'
        ozet_heading.runs[0].font.size = Pt(24)
        ozet_heading.runs[0].font.color.rgb = self.colors['primary']
        ozet_heading.space_after = Pt(20) # Başlık sonrası boşluk


        ozet_sections = {
            'Temel Değerlendirme': self.analiz_ozeti.get('temel_ozet', 'Değerlendirme bulunamadı.'),
            'Yatırım Değerlendirmesi': self.analiz_ozeti.get('yatirim_ozet', 'Değerlendirme bulunamadı.'),
            'Uygunluk Puanı Özeti': self.analiz_ozeti.get('uygunluk_ozet', 'Uygunluk puanı özeti bulunamadı.'), # Uygunluk puanı özeti eklendi
            'Öneriler ve Tavsiyeler': self.analiz_ozeti.get('tavsiyeler', 'Tavsiye bulunamadı.')
        }

        for title, text in ozet_sections.items():
            # Alt başlık
            sub_heading = doc.add_heading(title, level=2)
            run = sub_heading.runs[0]
            run.font.name = 'Arial'
            run.font.size = Pt(16)
            run.font.color.rgb = self.colors['secondary'] # Alt başlık rengi
            sub_heading.paragraph_format.space_before = Pt(15) # Önceki bölümle boşluk
            sub_heading.paragraph_format.space_after = Pt(8) # Başlık sonrası boşluk

            # İçerik paragrafı
            # Metni satır sonlarına göre böl ve her birini ayrı paragraf olarak ekle
            # Bu, Word'de daha iyi formatlama sağlayabilir.
            text_lines = str(text).splitlines()
            for line in text_lines:
                if line.strip(): # Boş satırları atla
                    p = doc.add_paragraph(line)
                    p.runs[0].font.name = 'Arial'
                    p.runs[0].font.size = Pt(11)
                    p.paragraph_format.line_spacing = 1.15 # Satır aralığı
                    p.space_after = Pt(2) # Paragraf sonrası boşluk (azaltıldı)
                else: # Boş satırsa, biraz daha fazla boşluk bırak
                    doc.add_paragraph().space_after = Pt(5)
            # Son paragraftan sonra genel bir boşluk
            doc.add_paragraph().space_after = Pt(10)


    def _add_word_photos_page(self, doc):
        """Word belgesine arsa fotoğrafları sayfası ekler."""
        foto_heading = doc.add_heading('Arsa Fotoğrafları', 1)
        foto_heading.runs[0].font.name = 'Arial'
        foto_heading.runs[0].font.size = Pt(24)
        foto_heading.runs[0].font.color.rgb = self.colors['primary']
        foto_heading.space_after = Pt(20) # Başlık sonrası boşluk

        uploaded_images = self._get_uploaded_images()
        logging.info(f"Word'e eklenecek resim sayısı: {len(uploaded_images)}")

        if not uploaded_images:
            p = doc.add_paragraph("Bu arsa için fotoğraf bulunmamaktadır.")
            p.runs[0].font.name = 'Arial'
            p.runs[0].font.italic = True
            p.space_after = Pt(10)
            logging.info("Eklenecek resim bulunamadı, uyarı mesajı eklendi.")
            return # Resim yoksa fonksiyondan çık

        # Resimleri 2x2 grid şeklinde ekle (yatay sayfaya daha uygun)
        section = doc.sections[0]
        available_width_cm = section.page_width.cm - section.left_margin.cm - section.right_margin.cm
        # Resimler için ayrılacak toplam alanın genişliği
        table_width_cm = available_width_cm * 0.95 # Sayfa genişliğinin %95'ini kullan
        max_img_width_cm = table_width_cm / 2 # Her sütun için maksimum genişlik
        max_img_height_cm = 8 # Makul bir yükseklik sınırı (Cm)

        num_images = len(uploaded_images)
        num_cols = 2
        # Her resim ve altındaki açıklama için 2 satır kullanacağız
        num_rows_per_image_row = 2 # Resim + Açıklama
        num_image_rows = (num_images + num_cols - 1) // num_cols # Resimlerin yerleşeceği satır sayısı
        total_rows = num_image_rows * num_rows_per_image_row # Toplam tablo satırı

        if total_rows == 0: # Hiç resim yoksa (yukarıda kontrol edildi ama yine de)
            return

        table = doc.add_table(rows=total_rows, cols=num_cols)
        table.alignment = WD_TABLE_ALIGNMENT.CENTER
        table.autofit = False
        table.allow_autofit = False

        # Sütun genişliklerini ayarla
        for col_idx in range(num_cols):
            table.columns[col_idx].width = Cm(max_img_width_cm)
            logging.info(f"Sütun {col_idx} genişliği ayarlandı: {max_img_width_cm} cm")


        img_idx = 0
        # Tablo satırlarında ilerle (her resim ve açıklaması için 2 satır)
        for r_offset in range(num_image_rows): # 0, 1, 2 ... (resim satır indeksi)
            img_row_idx = r_offset * num_rows_per_image_row # Resmin bulunduğu tablo satırı
            cap_row_idx = img_row_idx + 1 # Açıklamanın bulunduğu tablo satırı

            for c in range(num_cols):
                if img_idx < num_images:
                    img_path = uploaded_images[img_idx]
                    img_cell = table.cell(img_row_idx, c) # Resim hücresi
                    cap_cell = table.cell(cap_row_idx, c) # Açıklama hücresi

                    logging.info(f"Resim {img_idx+1} ekleniyor: {img_path}")

                    # Hücre dikey hizalama
                    img_cell.vertical_alignment = 1 # WD_ALIGN_VERTICAL.CENTER
                    cap_cell.vertical_alignment = 0 # WD_ALIGN_VERTICAL.TOP

                    # Resim ekleme paragrafı
                    p_img = img_cell.paragraphs[0]
                    p_img.clear()
                    p_img.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    run_img = p_img.add_run()

                    try:
                        # Resmi ekle ve boyutlandır
                        with Image.open(img_path) as img_pil:
                            width_px, height_px = img_pil.size
                            aspect_ratio = width_px / height_px if height_px > 0 else 1

                        # Boyut hesaplama: Cm cinsinden
                        target_width_val_cm = max_img_width_cm
                        target_height_val_cm = target_width_val_cm / aspect_ratio

                        if target_height_val_cm > max_img_height_cm:
                            target_height_val_cm = max_img_height_cm
                            target_width_val_cm = target_height_val_cm * aspect_ratio

                        run_img.add_picture(img_path, width=Cm(target_width_val_cm), height=Cm(target_height_val_cm))
                        logging.info(f"Resim {img_idx+1} başarıyla eklendi (G:{target_width_val_cm:.2f}cm, Y:{target_height_val_cm:.2f}cm).")

                        # Açıklama ekleme
                        p_cap = cap_cell.paragraphs[0]
                        p_cap.clear()
                        p_cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
                        run_cap = p_cap.add_run(f"Fotoğraf {img_idx + 1}")
                        run_cap.font.name = 'Arial'
                        run_cap.font.size = Pt(9)
                        run_cap.italic = True
                        p_cap.space_after = Pt(5) # Açıklama sonrası boşluk

                    except FileNotFoundError:
                        logging.error(f"Word resim dosyası bulunamadı: {img_path}")
                        # Hata mesajı ekle
                    except Exception as e:
                        logging.error(f"Word resim eklenemedi! {img_path}: {e}")
                        p_err = img_cell.paragraphs[0] # Hata mesajını resim hücresine yaz
                        p_err.clear()
                        p_err.alignment = WD_ALIGN_PARAGRAPH.CENTER
                        run_err = p_err.add_run(f"Resim Yüklenemedi\n({os.path.basename(img_path)})")
                        run_err.font.color.rgb = RGBColor(255, 0, 0)
                        run_err.font.size = Pt(9)
                        # Açıklama hücresini boş bırak veya hata mesajı ekle
                        p_cap_err = cap_cell.paragraphs[0]
                        p_cap_err.clear()
                        p_cap_err.alignment = WD_ALIGN_PARAGRAPH.CENTER
                        run_cap_err = p_cap_err.add_run("Yükleme Hatası")
                        run_cap_err.font.name = 'Arial'
                        run_cap_err.font.size = Pt(9)
                        run_cap_err.italic = True
                        run_cap_err.font.color.rgb = RGBColor(255, 0, 0)

                    img_idx += 1
                else: # Boş hücreler için (eğer resim sayısı tekse vs.)
                    table.cell(img_row_idx, c).text = ""
                    table.cell(cap_row_idx, c).text = ""


        # Tablo sonrası boşluk
        doc.add_paragraph().space_after = Pt(10)

    def _add_word_qr_code_page(self, doc, qr_path):
        """Word belgesine QR kod sayfası ekler."""
        heading = doc.add_heading('Analiz Detayları QR Kod', 1)
        heading.runs[0].font.name = 'Arial'
        heading.runs[0].font.size = Pt(24)
        heading.runs[0].font.color.rgb = self.colors['primary']
        heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
        heading.space_after = Pt(30) # Başlık sonrası boşluk

        # QR kod resmini ekle
        if os.path.exists(qr_path):
            try:
                p_qr = doc.add_paragraph()
                p_qr.alignment = WD_ALIGN_PARAGRAPH.CENTER
                run_qr = p_qr.add_run()
                run_qr.add_picture(qr_path, width=Inches(2.5)) # QR kod boyutunu ayarla
                p_qr.space_after = Pt(10) # QR kod sonrası boşluk
                logging.info("Word'e QR kod resmi eklendi.")
            except Exception as e:
                logging.error(f"Word QR kod eklenemedi: {e}")
                p_err = doc.add_paragraph("QR Kod Yüklenemedi.")
                p_err.alignment = WD_ALIGN_PARAGRAPH.CENTER
                p_err.runs[0].font.color.rgb = RGBColor(255, 0, 0)
        else:
             logging.warning("QR kod dosyası bulunamadı, Word'e eklenemedi.")
             p_err = doc.add_paragraph("QR Kod Yüklenemedi (Dosya Bulunamadı).")
             p_err.alignment = WD_ALIGN_PARAGRAPH.CENTER
             p_err.runs[0].font.color.rgb = RGBColor(255, 0, 0)

        # Açıklama ekle
        p_caption = doc.add_paragraph("Bu QR kodu tarayarak analiz detaylarına ulaşabilirsiniz.")
        p_caption.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p_caption.runs[0].font.name = 'Arial'
        p_caption.runs[0].font.size = Pt(11)
        p_caption.runs[0].font.italic = True
        p_caption.runs[0].font.color.rgb = RGBColor(120, 120, 120)


    def _register_fonts(self):
        """PDF için fontları kaydet - platform bağımsız (Türkçe karakterler için önemli)."""
        # Yaygın Türkçe destekli fontları dene
        fonts_to_try = {
            'Arial': ('arial.ttf', 'arialbd.ttf'),
            'DejaVuSans': ('DejaVuSans.ttf', 'DejaVuSans-Bold.ttf'), # Linux için yaygın
            'Calibri': ('calibri.ttf', 'calibrib.ttf'), # Windows için yaygın
            'Tahoma': ('tahoma.ttf', 'tahomabd.ttf'),
            'Verdana': ('verdana.ttf', 'verdanab.ttf'),
        }

        registered_base = None
        registered_bold = None

        font_paths = []
        # Sistem font dizinlerini ekle
        if 'SYSTEMROOT' in os.environ: # Windows
            font_paths.append(os.path.join(os.environ['SYSTEMROOT'], 'Fonts'))
        elif sys.platform == 'darwin': # macOS
            font_paths.append('/Library/Fonts')
            font_paths.append('/System/Library/Fonts')
            font_paths.append(os.path.expanduser('~/Library/Fonts'))
        else: # Linux ve diğerleri
            font_paths.append('/usr/share/fonts/truetype')
            font_paths.append('/usr/share/fonts/opentype')
            font_paths.append(os.path.expanduser('~/.fonts'))
            # Dağıtıma özel yollar eklenebilir (örn: /usr/local/share/fonts)

        logging.info(f"PDF fontları aranacak yollar: {font_paths}")

        for name, files in fonts_to_try.items():
            base_file, bold_file = files
            found_base = False
            found_bold = False
            base_path = None
            bold_path = None

            # Belirtilen yollarda ve alt klasörlerde fontları ara
            for path in font_paths:
                if not found_base:
                    search_pattern_base = os.path.join(path, '**', base_file)
                    possible_paths_base = glob.glob(search_pattern_base, recursive=True)
                    if possible_paths_base:
                        base_path = possible_paths_base[0]
                        try:
                            pdfmetrics.registerFont(TTFont(name, base_path))
                            registered_base = name
                            found_base = True
                            logging.info(f"PDF fontu bulundu ve kaydedildi (Normal): {name} -> {base_path}")
                        except Exception as e:
                            logging.warning(f"PDF fontu kaydedilemedi ({base_path}): {e}")
                
                if not found_bold:
                    search_pattern_bold = os.path.join(path, '**', bold_file)
                    possible_paths_bold = glob.glob(search_pattern_bold, recursive=True)
                    if possible_paths_bold:
                        bold_path = possible_paths_bold[0]
                        try:
                            pdfmetrics.registerFont(TTFont(name + '-Bold', bold_path)) # ReportLab bold suffix
                            registered_bold = name + '-Bold'
                            found_bold = True
                            logging.info(f"PDF fontu bulundu ve kaydedildi (Kalın): {name}-Bold -> {bold_path}")
                        except Exception as e:
                            logging.warning(f"PDF fontu kaydedilemedi ({bold_path}): {e}")

                # Hem normal hem kalın bulunduysa bu font setini kullan ve çık
                if found_base and found_bold:
                    logging.info(f"PDF için kullanılacak font seti: {registered_base}, {registered_bold}")
                    return registered_base, registered_bold

        # Hiç uygun set bulunamadıysa, varsayılana dön
        logging.warning("Sistemde uygun TTF font seti bulunamadı. Helvetica kullanılacak (Türkçe karakter sorunu olabilir).")
        return 'Helvetica', 'Helvetica-Bold'


    def create_pdf(self):
        """PDF belgesi oluşturur"""
        # Font ayarları
        print("\nDEBUG [Create PDF]: PDF belgesi oluşturma başladı.", flush=True) # Flush
        base_font, base_font_bold = self._register_fonts()

        from reportlab.platypus import BaseDocTemplate, PageTemplate, Frame

        filename = os.path.join(self.sunum_klasoru, f'analiz_{self.file_id}.pdf')
        doc = BaseDocTemplate(
            filename,
            pagesize= A4, # Dikey A4
            leftMargin=1.5*cm,
            rightMargin=1.5*cm,
            topMargin=2*cm,
            bottomMargin=2*cm,
            title=f"Gayrimenkul Analiz Raporu - {self.arsa_data.get('il', '')}/{self.arsa_data.get('ilce', '')}",
            author="İnveco Proje"
        )
        frame = Frame(doc.leftMargin, doc.bottomMargin, doc.width, doc.height, id='normal')
        doc.addPageTemplates([PageTemplate(id='footer', frames=frame, onPage=self._add_footer_pdf)])

        # Stil tanımlamaları (Türkçe karakterler için encoding='utf-8')
        styles = getSampleStyleSheet()
        styles.add(ParagraphStyle(
            name='CustomTitle',
            parent=styles['h1'],
            fontName=base_font_bold,
            fontSize=24,
            alignment=TA_CENTER,
            spaceAfter=20,
            textColor=colors.HexColor('#1A237E')
        ))

        styles.add(ParagraphStyle(
            name='CustomSubTitle',
            parent=styles['h2'],
            fontName=base_font,
            fontSize=18,
            alignment=TA_CENTER,
            spaceAfter=30,
            textColor=colors.HexColor('#3F51B5')
        ))

        styles.add(ParagraphStyle(
            name='CustomHeading1',
            parent=styles['h2'],
            fontName=base_font_bold,
            fontSize=16,
            spaceBefore=20,
            spaceAfter=10,
            textColor=colors.HexColor('#1A237E'),
            alignment=TA_LEFT,
            # encoding='utf-8' # ReportLab 4+ için gereksiz olabilir
        ))

        styles.add(ParagraphStyle(
            name='CustomHeading2',
            parent=styles['h3'],
            fontName=base_font_bold,
            fontSize=14,
            spaceBefore=15,
            spaceAfter=8,
            textColor=colors.darkslategray,
            alignment=TA_LEFT,
            # encoding='utf-8'
        ))

        styles.add(ParagraphStyle(
            name='CustomNormal',
            parent=styles['Normal'],
            fontName=base_font,
            fontSize=11,
            spaceAfter=6,
            # encoding='utf-8',
            leading=14, # Satır yüksekliği
            alignment=TA_LEFT
        ))

        styles.add(ParagraphStyle(
            name='CustomBullet',
            parent=styles['Bullet'],
            fontName=base_font,
            fontSize=11,
            spaceAfter=4,
            # encoding='utf-8',
            leading=14,
            leftIndent=20 # Madde işareti için girinti
        ))

        styles.add(ParagraphStyle(
            name='ImageCaption',
            parent=styles['Normal'],
            fontName=base_font,
            fontSize=9,
            alignment=TA_CENTER,
            spaceBefore=4,
            spaceAfter=10,
            textColor=colors.dimgray
        ))

        # Tablo stili (Arsa Bilgileri, Altyapı)
        info_table_style = TableStyle([
            # Başlık Satırı (Yok, direkt veri)
            # Veri Satırları
            ('BACKGROUND', (0, 0), (0, -1), colors.HexColor('#E8EAF6')),  # İlk sütun arka planı (Açık mavi)
            ('BACKGROUND', (1, 0), (1, -1), colors.white),              # İkinci sütun arka planı (Beyaz)
            ('TEXTCOLOR', (0, 0), (-1, -1), colors.black),
            ('ALIGN', (0, 0), (0, -1), 'LEFT'),                         # İlk sütun sola yaslı
            ('ALIGN', (1, 0), (1, -1), 'LEFT'),                         # İkinci sütun sola yaslı
            ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),                     # Dikey ortala
            ('FONTNAME', (0, 0), (0, -1), base_font_bold),              # İlk sütun kalın
            ('FONTNAME', (1, 0), (1, -1), base_font),                   # İkinci sütun normal
            ('FONTSIZE', (0, 0), (-1, -1), 11),
            ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#BDBDBD')), # İnce gri grid
            ('LEFTPADDING', (0, 0), (-1, -1), 8),
            ('RIGHTPADDING', (0, 0), (-1, -1), 8),
            ('TOPPADDING', (0, 0), (-1, -1), 6),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
        ])

        # Altyapı Tablosu için ek stil (Durum sütunu ortalı)
        altyapi_table_style = TableStyle([
            ('ALIGN', (1, 0), (1, -1), 'CENTER'), # Durum sütunu ortalı
        ], parent=info_table_style)


        elements = [] # PDF'e eklenecek akış öğeleri

        # Kapak Sayfası
        if self.logo_path and os.path.exists(self.logo_path):
            elements.append(RLImage(self.logo_path, width=4*cm, height=4*cm))
            elements.append(Spacer(1, 0.3*cm))
        elements.append(Paragraph('Arsa Analiz Raporu', styles['CustomTitle']))
        elements.append(Paragraph(f"{self.arsa_data.get('il', '')}, {self.arsa_data.get('ilce', '')}", styles['CustomSubTitle']))
        elements.append(Spacer(1, 0.5*cm))
        elements.append(Paragraph(datetime.now().strftime('%d.%m.%Y'), styles['CustomNormal']))
        elements.append(PageBreak())

        # PROFİL SAYFASI
        elements.append(Paragraph('Portföy Sorumlusu', styles['CustomHeading1']))
        # Profil fotoğrafı (varsa)
        profile_photo_path = self._get_profile_photo_path()
        if profile_photo_path:
            elements.append(RLImage(profile_photo_path, width=2.5*cm, height=2.5*cm))
            elements.append(Spacer(1, 0.2*cm))
        # Profil bilgileri tablosu
        profile_data = self._get_profile_table_data()
        table = Table(profile_data, colWidths=[doc.width*0.35, doc.width*0.65])
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (0, -1), colors.HexColor('#E3F2FD')),
            ('BACKGROUND', (1, 0), (1, -1), colors.white),
            ('TEXTCOLOR', (0, 0), (-1, -1), colors.black),
            ('ALIGN', (0, 0), (0, -1), 'LEFT'),
            ('ALIGN', (1, 0), (1, -1), 'LEFT'),
            ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
            ('FONTNAME', (0, 0), (0, -1), base_font_bold),
            ('FONTNAME', (1, 0), (1, -1), base_font),
            ('FONTSIZE', (0, 0), (-1, -1), 11),
            ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#BDBDBD')),
            ('LEFTPADDING', (0, 0), (-1, -1), 8),
            ('RIGHTPADDING', (0, 0), (-1, -1), 8),
            ('TOPPADDING', (0, 0), (-1, -1), 6),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
        ]))
        elements.append(table)
        elements.append(Spacer(1, 0.8*cm))
        elements.append(PageBreak())

        # Arsa Bilgileri Tablosu
        elements.append(Paragraph('Arsa Bilgileri', styles['CustomHeading1']))
        arsa_bilgileri_data = self._get_arsa_bilgileri() # [['Özellik', 'Değer'], ...]
        table = Table(arsa_bilgileri_data, colWidths=[doc.width*0.3, doc.width*0.7])
        table.setStyle(info_table_style)
        elements.append(table)
        elements.append(Spacer(1, 0.8*cm))

        # İnşaat Alanı Hesaplama Tablosu
        elements.append(Paragraph('İnşaat Alanı Hesaplaması', styles['CustomHeading1']))
        insaat_data = self._get_insaat_hesaplama()
        table = Table(insaat_data, colWidths=[doc.width*0.4, doc.width*0.6])
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (0, -1), colors.HexColor('#E3F2FD')),
            ('BACKGROUND', (1, 0), (1, -1), colors.white),
            ('TEXTCOLOR', (0, 0), (-1, -1), colors.black),
            ('ALIGN', (0, 0), (0, -1), 'LEFT'),
            ('ALIGN', (1, 0), (1, -1), 'LEFT'),
            ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
            ('FONTNAME', (0, 0), (0, -1), base_font_bold),
            ('FONTNAME', (1, 0), (1, -1), base_font),
            ('FONTSIZE', (0, 0), (-1, -1), 11),
            ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#BDBDBD')),
            ('LEFTPADDING', (0, 0), (-1, -1), 8),
            ('RIGHTPADDING', (0, 0), (-1, -1), 8),
            ('TOPPADDING', (0, 0), (-1, -1), 6),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
        ]))
        elements.append(table)
        elements.append(Spacer(1, 0.8*cm))

        # Altyapı Durumu Tablosu
        elements.append(Paragraph('Altyapı Durumu', styles['CustomHeading1']))
        altyapi_data = self._get_altyapi_durumu() # [['Özellik', 'Durum'], ...]
        # Durumları daha açıklayıcı yapalım
        altyapi_data_formatted = []
        for item, status in altyapi_data:
             status_text = "✓ Var" if status == '✓' else "✗ Yok"
             # Renklendirme için Paragraph kullanabiliriz ama tablo stilinde zor
             altyapi_data_formatted.append([item, status_text])

        table = Table(altyapi_data_formatted, colWidths=[doc.width*0.3, doc.width*0.7])
        table.setStyle(altyapi_table_style)
        elements.append(table)
        elements.append(Spacer(1, 0.8*cm))

        # SWOT Analizi - 2x2 tablo
        elements.append(PageBreak())
        elements.append(Paragraph('SWOT Analizi', styles['CustomHeading1']))

        swot_titles = ['Güçlü Yönler', 'Zayıf Yönler', 'Fırsatlar', 'Tehditler']
        swot_colors = [colors.HexColor('#A5D6A7'), colors.HexColor('#EF9A9A'),
                       colors.HexColor('#90CAF9'), colors.HexColor('#FFD54F')]
        swot_data = self._swot_table_data()
        swot_table_data = []
        for i in range(2):
            row = []
            for j in range(2):
                idx = i*2 + j
                title = f"<b>{swot_titles[idx]}</b><br/>"
                items = swot_data[i][j]
                if items:
                    text = title + "<br/>".join(f"• {madde}" for madde in items)
                else:
                    text = title + " - Yok -"
                row.append(Paragraph(text, styles['CustomNormal']))
            swot_table_data.append(row)
        swot_table = Table(swot_table_data, colWidths=[doc.width/2]*2, rowHeights=None)
        swot_table.setStyle(TableStyle([
            ('BACKGROUND', (0,0), (0,0), swot_colors[0]),
            ('BACKGROUND', (1,0), (1,0), swot_colors[1]),
            ('BACKGROUND', (0,1), (0,1), swot_colors[2]),
            ('BACKGROUND', (1,1), (1,1), swot_colors[3]),
            ('BOX', (0,0), (-1,-1), 1, colors.HexColor('#BDBDBD')),
            ('INNERGRID', (0,0), (-1,-1), 0.5, colors.HexColor('#BDBDBD')),
            ('VALIGN', (0,0), (-1,-1), 'TOP'),
            ('LEFTPADDING', (0,0), (-1,-1), 8),
            ('RIGHTPADDING', (0,0), (-1,-1), 8),
            ('TOPPADDING', (0,0), (-1,-1), 6),
            ('BOTTOMPADDING', (0,0), (-1,-1), 6),
        ]))
        elements.append(swot_table)
        elements.append(Spacer(1, 0.5*cm))

        # Analiz Özeti
        elements.append(PageBreak())
        elements.append(Paragraph('Analiz Özeti', styles['CustomHeading1']))

        ozet_sections_pdf = {
            'Temel Değerlendirme': self.analiz_ozeti.get('temel_ozet', 'Değerlendirme bulunamadı.'),
            'Yatırım Değerlendirmesi': self.analiz_ozeti.get('yatirim_ozet', 'Değerlendirme bulunamadı.'),
            'Öneriler ve Tavsiyeler': self.analiz_ozeti.get('tavsiyeler', 'Tavsiye bulunamadı.')
        }
        for title, text in ozet_sections_pdf.items():
             elements.append(Paragraph(title, styles['CustomHeading2']))
             elements.append(Paragraph(text.replace('\n', '<br/>'), styles['CustomNormal'])) # Satır sonlarını <br/> ile değiştir
             elements.append(Spacer(1, 0.4*cm))


        # Arsa Fotoğrafları Bölümü
        elements.append(PageBreak())
        elements.append(Paragraph('Gayrimenkul Fotoğrafları', styles['CustomHeading1']))

        # --- YENİ LOG ---
        print("DEBUG [Create PDF]: _get_uploaded_images çağrılıyor...")
        uploaded_images = self._get_uploaded_images() # Kopyalanmış resim yolları
        # --- YENİ LOG ---
        print(f"DEBUG [Create PDF]: _get_uploaded_images'dan dönen resimler: {uploaded_images}", flush=True)

        if not uploaded_images:
            elements.append(Paragraph("Bu arsa için fotoğraf bulunmamaktadır.", styles['CustomNormal']))
            # --- YENİ LOG ---
            print("DEBUG [Create PDF]: Eklenecek resim bulunamadı.")
        else:
            # --- YENİ LOG ---
            print(f"DEBUG [Create PDF]: {len(uploaded_images)} adet resim eklenecek.")
            # Resimleri ekle
            max_img_width = doc.width * 0.9
            max_img_height = doc.height * 0.4


            for i, img_path in enumerate(uploaded_images):
                 # --- YENİ LOG ---
                print(f"DEBUG [Create PDF]: Resim {i+1} ekleniyor: {img_path}")
                print(f"DEBUG [Create PDF]: Bu dosya var mı? {os.path.exists(img_path)}")
                try:
                    # --- YENİ LOG ---
                    print(f"DEBUG [Create PDF]: Pillow ile açılıyor: {img_path}")
                    with Image.open(img_path) as pil_img:
                        width_px, height_px = pil_img.size
                        aspect_ratio = width_px / height_px
                        print(f"DEBUG [Create PDF]: Orijinal boyut (px): {width_px}x{height_px}, Oran: {aspect_ratio:.2f}")

                    # Boyutlandırma
                    img_width = max_img_width
                    img_height = img_width / aspect_ratio
                    if img_height > max_img_height:
                        img_height = max_img_height
                        img_width = img_height * aspect_ratio
                    if img_width > max_img_width:
                         img_width = max_img_width
                         img_height = img_width / aspect_ratio

                    # --- YENİ LOG ---
                    print(f"DEBUG [Create PDF]: Hedef boyut (pt): {img_width:.2f}x{img_height:.2f}")
                    print(f"DEBUG [Create PDF]: RLImage oluşturuluyor...")
                    rl_image = RLImage(img_path, width=img_width, height=img_height)
                    elements.append(rl_image)
                    # --- YENİ LOG ---
                    print(f"DEBUG [Create PDF]: Resim {i+1} başarıyla elementlere eklendi.")

                    # Açıklama
                    caption_text = f"Fotoğraf {i+1}: {os.path.splitext(os.path.basename(img_path))[0]}"
                    elements.append(Paragraph(caption_text, styles['ImageCaption']))
                    elements.append(Spacer(1, 0.5*cm))

                    # Her 2 resimde bir sayfa sonu ekle (isteğe bağlı)
                    # if (i + 1) % 2 == 0 and i + 1 < len(uploaded_images):
                    #     elements.append(PageBreak())

                except Exception as e:
                    print(f"HATA [PDF Resim Ekleme]: {img_path} eklenemedi: {e}")
                    error_text = f"Resim Yüklenemedi ({os.path.basename(img_path)}) - Hata: {e}"
                    elements.append(Paragraph(error_text, styles['CustomNormal']))

        # PDF'i oluştur
        # PDF'i oluştur
        # --- YENİ LOG ---
        print(f"DEBUG [Create PDF]: PDF dosyası build ediliyor: {filename}", flush=True)
        try:
            doc.build(elements)
            print(f"PDF dosyası başarıyla kaydedildi: {filename}", flush=True)
        except Exception as e:
            print(f"HATA [PDF Build]: PDF oluşturulamadı: {e}", flush=True)
            traceback.print_exc(file=sys.stdout)
            sys.stdout.flush()
            raise

        print("DEBUG [Create PDF]: create_pdf metodu tamamlandı.", flush=True)

        return filename
    def _add_profile_section_pdf(self, styles):
        """PDF belgesine profil bilgilerini ekler"""
        elements = []
        styleHeading1 = styles['Heading1']
        styleN = styles['Normal']
        styleN.fontName = 'DejaVuSans' # Fontu ayarla

        elements.append(Paragraph("Analist Bilgileri", styleHeading1))

        # Profil fotoğrafı ekle (varsa)
        profile_photo_path = self._get_profile_photo_path()
        if profile_photo_path and os.path.exists(profile_photo_path):
            try:
                img = RLImage(profile_photo_path, width=1.5*inch, height=1.5*inch) # Boyutu ayarla
                img.hAlign = 'CENTER'
                elements.append(img)
                elements.append(Spacer(1, 0.2*inch))
                print(f"DEBUG [Add Profile PDF]: Profil fotoğrafı eklendi: {profile_photo_path}")
            except Exception as e:
                print(f"HATA [Add Profile PDF]: Profil fotoğrafı eklenemedi: {e}")
                traceback.print_exc(file=sys.stdout) # Konsola yazdır
                sys.stdout.flush() # Zorla


        data = self._get_profile_table_data()
        table_data = [['Özellik', 'Değer']] # Başlık satırı
        for label, value in data:
            table_data.append([label, value])

        table = Table(table_data, colWidths=[4*cm, 8*cm]) # Sütun genişliklerini ayarla
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor(self.colors['background'])), # Başlık satırı arkaplan
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.HexColor('#475569')), # Başlık satırı yazı rengi
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (-1, 0), 'DejaVuSans-Bold'), # Başlık fontu
            ('FONTNAME', (0, 1), (-1, -1), 'DejaVuSans'), # İçerik fontu
            ('FONTSIZE', (0, 0), (-1, -1), 10),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
            ('TOPPADDING', (0, 0), (-1, -1), 6),
            ('BACKGROUND', (0, 1), (-1, -1), colors.white), # İçerik satırları arkaplan
            ('GRID', (0, 0), (-1, -1), 1, colors.black), # Tablo çizgileri
        ]))
        elements.append(table)
        elements.append(Spacer(1, 0.5*inch)) # Tablodan sonra boşluk
        return elements

    def _add_property_section_pdf(self, styles):
        """PDF belgesine arsa bilgilerini ekler"""
        elements = []
        styleHeading1 = styles['Heading1']
        styleN = styles['Normal']
        styleN.fontName = 'DejaVuSans' # Fontu ayarla

        elements.append(Paragraph("Arsa Bilgileri", styleHeading1))

        data = self._get_arsa_bilgileri()
        table_data = [['Özellik', 'Değer']] # Başlık satırı
        for label, value in data:
            table_data.append([label, value])

        table = Table(table_data, colWidths=[4*cm, 8*cm]) # Sütun genişliklerini ayarla
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor(self.colors['background'])),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.HexColor('#475569')),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (-1, 0), 'DejaVuSans-Bold'),
            ('FONTNAME', (0, 1), (-1, -1), 'DejaVuSans'),
            ('FONTSIZE', (0, 0), (-1, -1), 10),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
            ('TOPPADDING', (0, 0), (-1, -1), 6),
            ('BACKGROUND', (0, 1), (-1, -1), colors.white),
            ('GRID', (0, 0), (-1, -1), 1, colors.black),
        ]))
        elements.append(table)
        elements.append(Spacer(1, 0.5*inch))
        return elements

    def _add_infrastructure_section_pdf(self, styles):
        """PDF belgesine altyapı durumunu ekler"""
        elements = []
        styleHeading1 = styles['Heading1']
        styleN = styles['Normal']
        styleN.fontName = 'DejaVuSans' # Fontu ayarla

        elements.append(Paragraph("Altyapı Durumu", styleHeading1))

        data = self._get_altyapi_durumu()
        table_data = [['Altyapı', 'Durum']] # Başlık satırı
        for label, status in data:
            table_data.append([label, status])

        table = Table(table_data, colWidths=[6*cm, 6*cm]) # Sütun genişliklerini ayarla
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor(self.colors['background'])),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.HexColor('#475569')),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (-1, 0), 'DejaVuSans-Bold'),
            ('FONTNAME', (0, 1), (-1, -1), 'DejaVuSans'),
            ('FONTSIZE', (0, 0), (-1, -1), 10),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
            ('TOPPADDING', (0, 0), (-1, -1), 6),
            ('BACKGROUND', (0, 1), (-1, -1), colors.white),
            ('GRID', (0, 0), (-1, -1), 1, colors.black),
            # Duruma göre hücre rengi (PDF'te hücre arkaplanı renklemek daha karmaşık, şimdilik sadece yazı rengi)
            ('TEXTCOLOR', (1, 1), (1, -1), colors.green), # ✓ için yeşil
            ('TEXTCOLOR', (1, 1), (1, -1), colors.red), # ✗ için kırmızı (Bu şekilde olmaz, her satır için ayrı stil gerekir)
            # Alternatif olarak, her hücre için ayrı stil tanımlanabilir veya sadece metin rengi değiştirilebilir.
            # Şimdilik sadece metin rengini değiştirelim.
        ]))

        # Her hücrenin metin rengini duruma göre ayarla
        for i in range(1, len(table_data)):
             status = table_data[i][1]
             if status == '✓':
                 table.setStyle(TableStyle([('TEXTCOLOR', (1, i), (1, i), colors.green)]))
             else:
                 table.setStyle(TableStyle([('TEXTCOLOR', (1, i), (1, i), colors.red)]))


        elements.append(table)
        elements.append(Spacer(1, 0.5*inch))
        return elements

    def _add_swot_section_pdf(self, styles):
        """PDF belgesine SWOT analizini ekler"""
        elements = []
        styleHeading1 = styles['Heading1']
        styleHeading2 = styles['Heading2']
        styleN = styles['Normal']
        styleN.fontName = 'DejaVuSans' # Fontu ayarla
        styleHeading2.fontName = 'DejaVuSans-Bold' # Başlık fontu

        elements.append(Paragraph("SWOT Analizi", styleHeading1))

        swot_data = self._get_swot_analizi()

        # SWOT başlıkları ve renkleri
        swot_titles = {
            'Güçlü Yönler': colors.green,
            'Zayıf Yönler': colors.red,
            'Fırsatlar': colors.blue,
            'Tehditler': colors.orange
        }

        for title, color in swot_titles.items():
            elements.append(Paragraph(title, ParagraphStyle(
                name=f'{title.replace(" ", "")}Style',
                parent=styleHeading2,
                textColor=color
            )))
            items = swot_data.get(title, [])
            if items:
                for item in items:
                    elements.append(Paragraph(f"- {item}", styleN))
            else:
                elements.append(Paragraph("Veri bulunamadı.", styleN))
            elements.append(Spacer(1, 0.2*inch)) # Maddeler arası boşluk

        elements.append(Spacer(1, 0.5*inch)) # Bölüm sonu boşluk
        return elements

    def _add_photos_section_pdf(self, styles):
        """PDF belgesine yüklenen fotoğrafları ekler"""
        print("DEBUG [Add Photos PDF]: _add_photos_section_pdf metodu başladı.")
        elements = []
        styleHeading1 = styles['Heading1']
        styleN = styles['Normal']
        styleN.fontName = 'DejaVuSans' # Fontu ayarla

        elements.append(Paragraph("Arsa Fotoğrafları", styleHeading1))

        images = self._get_uploaded_images()
        print(f"DEBUG [Add Photos PDF]: Eklenecek resim sayısı: {len(images)}") # Flush

        if not images:
            elements.append(Paragraph("Yüklü fotoğraf bulunamadı.", styleN))
            print("DEBUG [Add Photos PDF]: Yüklü fotoğraf bulunamadı mesajı eklendi.") # Flush
            return elements

        # Resimleri 2 sütunlu bir tabloya yerleştirme
        table_data = []
        row = []
        for i in range(0, len(images), 2):
            row_elements = []
            # İlk sütun
            img_path1 = images[i]
            try:
                # Resim boyutunu ayarla (örn: 8 cm genişlik)
                img1 = RLImage(img_path1, width=8*cm, height=8*cm, kind='proportional')
                img1.hAlign = 'CENTER'
                row_elements.append(img1)
                print(f"DEBUG [Add Photos PDF]: Resim 1 eklendi: {img_path1}") # Flush
            except Exception as e:
                print(f"HATA [Add Photos PDF]: Resim 1 eklenemedi: {img_path1} - {e}") # Flush
                traceback.print_exc(file=sys.stdout) # Konsola yazdır
                sys.stdout.flush() # Zorla
                row_elements.append(Paragraph("Resim yüklenemedi.", styleN)) # Hata mesajı ekle

            # İkinci sütun (varsa)
            if i + 1 < len(images):
                img_path2 = images[i+1]
                try:
                    img2 = RLImage(img_path2, width=8*cm, height=8*cm, kind='proportional')
                    img2.hAlign = 'CENTER'
                    row_elements.append(img2)
                    print(f"DEBUG [Add Photos PDF]: Resim 2 eklendi: {img_path2}") # Flush
                except Exception as e:
                    print(f"HATA [Add Photos PDF]: Resim 2 eklenemedi: {img_path2} - {e}") # Flush
                    traceback.print_exc(file=sys.stdout) # Konsola yazdır
                    sys.stdout.flush() # Zorla
                    row_elements.append(Paragraph("Resim yüklenemedi.", styleN)) # Hata mesajı ekle
            else:
                 # Tek resim varsa ikinci hücre boş kalır
                 row_elements.append(Spacer(1,1)) # Boşluk ekle

            table_data.append(row_elements)


        if table_data:
            table = Table(table_data, colWidths=[9*cm, 9*cm]) # Sütun genişliklerini ayarla
            table.setStyle(TableStyle([
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                ('LEFTPADDING', (0, 0), (-1, -1), 0),
                ('RIGHTPADDING', (0, 0), (-1, -1), 0),
                ('TOPPADDING', (0, 0), (-1, -1), 0),
                ('BOTTOMPADDING', (0, 0), (-1, -1), 0),
                ('BACKGROUND', (0, 0), (-1, -1), colors.white), # Arkaplanı beyaz yap
            ]))
            elements.append(table)
            elements.append(Spacer(1, 0.5*inch)) # Tablodan sonra boşluk

        print("DEBUG [Add Photos PDF]: _add_photos_section_pdf metodu tamamlandı.") # Flush
        return elements

    def create_pptx(self):
        """PowerPoint sunumu oluşturur."""
        logging.info("PowerPoint sunumu oluşturma başladı.")
        try:
            # Mevcut sunum şablonunu yükle
            # Şablon yolu app.py'nin bulunduğu dizine göre ayarlanmalı
            base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
            template_path = os.path.join(base_dir, 'templates', 'sunum', 'ppt_template.pptx')
            logging.info(f"PPTX şablon yolu: {template_path}")

            if not os.path.exists(template_path):
                logging.error(f"PPTX şablon dosyası bulunamadı: {template_path}")
                # Şablon yoksa boş bir sunum oluşturmayı deneyebiliriz veya None dönebiliriz.
                # Şimdilik None dönelim.
                # prs = Presentation() # Boş sunum oluştur
                return None


            prs = Presentation(template_path)

            # İlk slayt (Başlık Slaytı)
            # Şablondaki layout'ları kontrol etmek gerekebilir. Genellikle 0 indexindedir.
            try:
                title_slide_layout = prs.slide_layouts[0] # Başlık Slaytı
            except IndexError:
                logging.warning("Şablonda 0 indexli layout (Başlık Slaytı) bulunamadı, ilk uygun layout kullanılıyor.")
                title_slide_layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0]


            slide = prs.slides.add_slide(title_slide_layout)

            # Başlık ve Alt Başlık Placeholder'larını bul ve doldur
            try:
                title_shape = slide.shapes.title
                if title_shape:
                    title_shape.text = "GAYRİMENKUL ANALİZ RAPORU"
                    logging.info("PPTX başlık placeholder dolduruldu.")
                else: # title placeholder yoksa
                    logging.warning("PPTX başlık placeholder (title) bulunamadı. Manuel ekleniyor.")
                    left, top, width, height = PPTXInches(1), PPTXInches(2), PPTXInches(8), PPTXInches(1.5)
                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    tf.clear() # Önceki içeriği temizle (gerekirse)
                    p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
                    p.text = "GAYRİMENKUL ANALİZ RAPORU"
                    p.font.size = PPTXPt(44)
                    p.font.bold = True
                    p.alignment = WD_ALIGN_PARAGRAPH.CENTER # python-pptx'te bu enum yok, int değeri kullanılmalı: 1 (LEFT), 2 (CENTER), 3 (RIGHT)
                    p.alignment = 2 # CENTER
            except AttributeError: # shapes.title yoksa
                logging.warning("PPTX başlık placeholder (shapes.title) bulunamadı. Manuel ekleniyor.")
                left, top, width, height = PPTXInches(1), PPTXInches(2), PPTXInches(8), PPTXInches(1.5)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                p = tf.add_paragraph()
                p.text = "GAYRİMENKUL ANALİZ RAPORU"
                p.font.size = PPTXPt(44); p.font.bold = True; p.alignment = 2


            # Alt başlık placeholder'ı
            # Genellikle ilk placeholder (title olmayan) alt başlık olur.
            subtitle_placeholder = None
            if len(slide.placeholders) > 1:
                for i in range(len(slide.placeholders)):
                    # Placeholder'ın title olmadığını ve metin kutusu olduğunu kontrol et
                    # (slide.shapes.title zaten title_shape ile alındı)
                    # Placeholder ID'leri şablona göre değişir.
                    # Genellikle title'dan sonraki ilk placeholder alt başlık olabilir.
                    if slide.placeholders[i] != title_shape and hasattr(slide.placeholders[i], 'text_frame'):
                        subtitle_placeholder = slide.placeholders[i]
                        break # İlk uygun olanı al

            if subtitle_placeholder:
                tf = subtitle_placeholder.text_frame
                tf.clear()
                p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
                p.text = f"{self.arsa_data.get('il', '')}, {self.arsa_data.get('ilce', '')}"
                p.font.size = PPTXPt(28)
                p.alignment = 2 # CENTER

                p2 = tf.add_paragraph()
                p2.text = datetime.now().strftime('%d.%m.%Y')
                p2.font.size = PPTXPt(20)
                p2.alignment = 2 # CENTER
                logging.info("PPTX alt başlık placeholder dolduruldu.")
            else:
                logging.warning("PPTX alt başlık placeholder bulunamadı. Manuel ekleniyor.")
                left, top, width, height = PPTXInches(1), PPTXInches(3.8), PPTXInches(8), PPTXInches(1.5)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                p = tf.add_paragraph()
                p.text = f"{self.arsa_data.get('il', '')}, {self.arsa_data.get('ilce', '')}"
                p.font.size = PPTXPt(28); p.alignment = 2
                p2 = tf.add_paragraph()
                p2.text = datetime.now().strftime('%d.%m.%Y')
                p2.font.size = PPTXPt(20); p2.alignment = 2


            # Logo ekle (varsa)
            if os.path.exists(self.logo_path):
                 try:
                     left = PPTXInches(0.5); top = PPTXInches(0.5)
                     slide.shapes.add_picture(self.logo_path, left, top, width=PPTXInches(1.2))
                     logging.info("PPTX kapak slaytına logo eklendi.")
                 except Exception as e:
                     logging.error(f"PPTX kapak slaytına logo eklenemedi: {e}")


            # PROFİL SLAYTI
            self._add_pptx_profile_slide(prs)

            # ARSA BİLGİLERİ SLAYTI
            self._add_pptx_property_slide(prs)

            # ALTYAPI DURUMU SLAYTI
            self._add_pptx_infrastructure_slide(prs)

            # İNŞAAT ALANI HESAPLAMA SLAYTI
            self._add_pptx_insaat_hesaplama_slide(prs)

            # SWOT ANALİZİ SLAYTI
            self._add_pptx_swot_slide(prs)

            # ANALİZ ÖZETİ SLAYTI
            self._add_pptx_analiz_ozeti_slide(prs)

            # ARSA FOTOĞRAFLARI SLAYTI
            self._add_pptx_photos_slide(prs)

            # QR KOD SLAYTI
            qr_path = self._create_qr_code()
            if qr_path and os.path.exists(qr_path):
                 self._add_pptx_qr_code_slide(prs, qr_path)


            # Dosya adını oluştur
            il = self.arsa_data.get('il', 'Bilinmiyor')
            ilce = self.arsa_data.get('ilce', 'Bilinmiyor')
            # Dosya adında özel karakter olmamasına dikkat et
            safe_il = "".join([c for c in il if c.isalnum() or c in (' ', '-')]).replace(' ', '_')
            safe_ilce = "".join([c for c in ilce if c.isalnum() or c in (' ', '-')]).replace(' ', '_')
            dosya_adi = f"arsa_analiz_{safe_il}_{safe_ilce}_{self.file_id}.pptx"
            filepath = os.path.join(self.sunum_klasoru, dosya_adi)

            # Sunumu kaydet
            prs.save(filepath)
            logging.info(f"PowerPoint sunumu kaydedildi: {filepath}")
            return filepath

        except Exception as e:
            logging.error(f"PowerPoint sunumu oluşturulurken genel hata: {e}")
            traceback.print_exc(file=sys.stdout)
            sys.stdout.flush()
            return None

    def _add_pptx_slide_with_title(self, prs, title_text, layout_index=1):
        """PPTX'e başlık içeren yeni bir slayt ekler ve slayt nesnesini döndürür."""
        # Başlık ve İçerik layout'u (genellikle 1 indexindedir, şablona göre değişebilir)
        try:
            slide_layout = prs.slide_layouts[layout_index]
        except IndexError:
            logging.warning(f"Şablonda {layout_index} indexli layout bulunamadı, ilk uygun (genellikle boş) layout kullanılıyor.")
            # Genellikle 5 veya 6 boş layout olur.
            fallback_layout_index = 5 if len(prs.slide_layouts) > 5 else (len(prs.slide_layouts) -1 if len(prs.slide_layouts) > 0 else 0)
            try:
                slide_layout = prs.slide_layouts[fallback_layout_index]
            except IndexError: # Hiç layout yoksa (çok olası değil ama)
                logging.error("Şablonda hiç layout bulunamadı!")
                # Bu durumda yeni slayt ekleyemeyiz, None dönebiliriz veya hata fırlatabiliriz.
                # Şimdilik None dönelim, çağıran fonksiyon kontrol etmeli.
                return None


        slide = prs.slides.add_slide(slide_layout)

        # Başlık placeholder'ını bul ve doldur
        try:
            title_shape = slide.shapes.title
            if title_shape:
                title_shape.text = title_text
            else: # title placeholder yoksa
                logging.warning(f"Slayt başlık placeholder (title) bulunamadı: '{title_text}'. Manuel ekleniyor.")
                left, top, width, height = PPTXInches(0.5), PPTXInches(0.3), PPTXInches(9), PPTXInches(0.9)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                p = tf.add_paragraph()
                p.text = title_text
                p.font.size = PPTXPt(32); p.font.bold = True; p.alignment = 2 # CENTER
        except AttributeError: # shapes.title yoksa
            logging.warning(f"Slayt başlık placeholder (shapes.title) bulunamadı: '{title_text}'. Manuel ekleniyor.")
            left, top, width, height = PPTXInches(0.5), PPTXInches(0.3), PPTXInches(9), PPTXInches(0.9)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            p = tf.add_paragraph()
            p.text = title_text
            p.font.size = PPTXPt(32); p.font.bold = True; p.alignment = 2 # CENTER

        return slide

    def _add_pptx_profile_slide(self, prs):
        """PowerPoint sunumuna profil bilgilerini içeren slayt ekler."""
        slide = self._add_pptx_slide_with_title(prs, "Portföy Sorumlusu")
        if not slide: return # Slayt oluşturulamadıysa çık

        # Profil fotoğrafı ekle (varsa) - Genellikle sağ üst veya sağ orta
        profile_photo_path = self._get_profile_photo_path()
        img_left, img_top, img_width, img_height = PPTXInches(7.2), PPTXInches(1.5), PPTXInches(2.3), PPTXInches(2.3) # Sağ taraf için konum
        if profile_photo_path:
            try:
                slide.shapes.add_picture(profile_photo_path, img_left, img_top, width=img_width, height=img_height)
                logging.info(f"PPTX profil fotoğrafı eklendi: {profile_photo_path}")
            except Exception as e:
                logging.error(f"PPTX profil fotoğrafı eklenemedi: {e}")


        # Tablo ekle - Fotoğrafın soluna
        data = self._get_profile_table_data()
        rows = len(data)
        cols = 2
        # Tablo konumu ve boyutu (Başlık altından başla, fotoğrafın solunda)
        table_left, table_top, table_width, table_height = PPTXInches(0.5), PPTXInches(1.5), PPTXInches(6.5), PPTXInches(len(data) * 0.5)


        try:
            shape = slide.shapes.add_table(rows, cols, table_left, table_top, table_width, table_height)
            table = shape.table

            # Sütun genişliklerini ayarla
            table.columns[0].width = PPTXInches(2.0)
            table.columns[1].width = PPTXInches(4.5)


            # Tablo verilerini doldur
            for i, (label, value) in enumerate(data):
                cell1 = table.cell(i, 0)
                cell2 = table.cell(i, 1)
                
                tf1 = cell1.text_frame; tf1.clear(); p1 = tf1.add_paragraph()
                p1.text = label
                p1.font.bold = True; p1.font.size = PPTXPt(11)
                
                tf2 = cell2.text_frame; tf2.clear(); p2 = tf2.add_paragraph()
                p2.text = str(value) # Değeri stringe çevir
                p2.font.size = PPTXPt(11)

                # Hücre dikey hizalamayı ayarla (MSO_ANCHOR enum değerleri)
                # 1: MIDDLE, 0: TOP, 2: BOTTOM
                cell1.vertical_anchor = 1 # MIDDLE
                cell2.vertical_anchor = 1 # MIDDLE

            # Tablo stilini ayarla (isteğe bağlı)
            # table.style = prs.table_styles[0] # Şablondaki ilk stili uygula (deneyerek bulunabilir)

        except Exception as e:
             logging.error(f"PPTX profil tablosu eklenemedi: {e}")
             txBox = slide.shapes.add_textbox(table_left, table_top, table_width, PPTXInches(2))
             tf = txBox.text_frame
             p = tf.add_paragraph(); p.text = "Profil bilgileri yüklenemedi."; p.font.size = PPTXPt(10)


    def _add_pptx_property_slide(self, prs):
        """PowerPoint sunumuna arsa bilgilerini içeren slayt ekler."""
        slide = self._add_pptx_slide_with_title(prs, "Arsa Bilgileri")
        if not slide: return

        data = self._get_arsa_bilgileri()
        rows = len(data)
        cols = 2
        # Tablo konumu ve boyutu (Slaytın ortasına yakın)
        table_left, table_top, table_width, table_height = PPTXInches(1), PPTXInches(1.5), PPTXInches(8), PPTXInches(len(data) * 0.55)


        try:
            shape = slide.shapes.add_table(rows, cols, table_left, table_top, table_width, table_height)
            table = shape.table

            table.columns[0].width = PPTXInches(2.5)
            table.columns[1].width = PPTXInches(5.5)

            for i, (label, value) in enumerate(data):
                cell1 = table.cell(i, 0); tf1 = cell1.text_frame; tf1.clear(); p1 = tf1.add_paragraph()
                p1.text = label; p1.font.bold = True; p1.font.size = PPTXPt(11)
                
                cell2 = table.cell(i, 1); tf2 = cell2.text_frame; tf2.clear(); p2 = tf2.add_paragraph()
                p2.text = str(value); p2.font.size = PPTXPt(11)

                cell1.vertical_anchor = 1; cell2.vertical_anchor = 1

        except Exception as e:
            logging.error(f"PPTX arsa bilgileri tablosu eklenemedi: {e}")
            txBox = slide.shapes.add_textbox(table_left, table_top, table_width, PPTXInches(2))
            tf = txBox.text_frame; p = tf.add_paragraph(); p.text = "Arsa bilgileri yüklenemedi."; p.font.size = PPTXPt(10)


    def _add_pptx_infrastructure_slide(self, prs):
        """PowerPoint sunumuna altyapı durumunu içeren slayt ekler."""
        slide = self._add_pptx_slide_with_title(prs, "Altyapı Durumu")
        if not slide: return

        data = self._get_altyapi_durumu()
        rows = len(data)
        cols = 2
        table_left, table_top, table_width, table_height = PPTXInches(1.5), PPTXInches(1.8), PPTXInches(7), PPTXInches(len(data) * 0.5)

        try:
            shape = slide.shapes.add_table(rows, cols, table_left, table_top, table_width, table_height)
            table = shape.table

            table.columns[0].width = PPTXInches(3.0)
            table.columns[1].width = PPTXInches(4.0)

            for i, (label, status) in enumerate(data):
                cell1 = table.cell(i, 0); tf1 = cell1.text_frame; tf1.clear(); p1 = tf1.add_paragraph()
                p1.text = label; p1.font.bold = True; p1.font.size = PPTXPt(11)
                
                cell2 = table.cell(i, 1); tf2 = cell2.text_frame; tf2.clear(); p2 = tf2.add_paragraph()
                status_text = '✓ Var' if status == '✓' else '✗ Yok'
                p2.text = status_text; p2.font.size = PPTXPt(11); p2.font.bold = True
                
                # Duruma göre metin rengini ayarla
                if status == '✓':
                     p2.font.color.rgb = RGBColor(0, 128, 0) # Yeşil
                else:
                     p2.font.color.rgb = RGBColor(255, 0, 0) # Kırmızı
                
                # Metni ortala
                p2.alignment = 2 # CENTER


                cell1.vertical_anchor = 1; cell2.vertical_anchor = 1

        except Exception as e:
            logging.error(f"PPTX altyapı tablosu eklenemedi: {e}")
            txBox = slide.shapes.add_textbox(table_left, table_top, table_width, PPTXInches(2))
            tf = txBox.text_frame; p = tf.add_paragraph(); p.text = "Altyapı bilgileri yüklenemedi."; p.font.size = PPTXPt(10)


    def _add_pptx_insaat_hesaplama_slide(self, prs):
        """PowerPoint sunumuna inşaat alanı hesaplamasını içeren slayt ekler."""
        slide = self._add_pptx_slide_with_title(prs, "İnşaat Alanı Hesaplaması")
        if not slide: return

        data = self._get_insaat_hesaplama()
        rows = len(data)
        cols = 2
        table_left, table_top, table_width, table_height = PPTXInches(1.5), PPTXInches(1.8), PPTXInches(7), PPTXInches(len(data) * 0.55)

        try:
            shape = slide.shapes.add_table(rows, cols, table_left, table_top, table_width, table_height)
            table = shape.table

            table.columns[0].width = PPTXInches(4.0) # Etiket için daha geniş
            table.columns[1].width = PPTXInches(3.0)

            for i, (label, value) in enumerate(data):
                cell1 = table.cell(i, 0); tf1 = cell1.text_frame; tf1.clear(); p1 = tf1.add_paragraph()
                p1.text = label; p1.font.bold = True; p1.font.size = PPTXPt(11)
                
                cell2 = table.cell(i, 1); tf2 = cell2.text_frame; tf2.clear(); p2 = tf2.add_paragraph()
                p2.text = str(value); p2.font.size = PPTXPt(11)

                cell1.vertical_anchor = 1; cell2.vertical_anchor = 1

        except Exception as e:
            logging.error(f"PPTX inşaat hesaplama tablosu eklenemedi: {e}")
            txBox = slide.shapes.add_textbox(table_left, table_top, table_width, PPTXInches(2))
            tf = txBox.text_frame; p = tf.add_paragraph(); p.text = "İnşaat hesaplama bilgileri yüklenemedi."; p.font.size = PPTXPt(10)


    def _add_pptx_swot_slide(self, prs):
        """PowerPoint sunumuna SWOT analizini içeren slayt ekler."""
        slide = self._add_pptx_slide_with_title(prs, "SWOT Analizi")
        if not slide: return

        swot_data_dict = self._get_swot_analizi()

        # SWOT başlıkları ve renkleri (RGBColor nesneleri)
        swot_config = {
            'Güçlü Yönler': {'color': RGBColor(0, 128, 0), 'key': 'Güçlü Yönler'}, # Yeşil
            'Zayıf Yönler': {'color': RGBColor(255, 0, 0), 'key': 'Zayıf Yönler'}, # Kırmızı
            'Fırsatlar': {'color': RGBColor(0, 0, 255), 'key': 'Fırsatlar'},   # Mavi
            'Tehditler': {'color': RGBColor(255, 165, 0), 'key': 'Tehditler'}    # Turuncu
        }
        swot_keys_ordered = ['Güçlü Yönler', 'Zayıf Yönler', 'Fırsatlar', 'Tehditler']


        # Slayt genişliği ve yüksekliği (yaklaşık değerler, şablona göre değişebilir)
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        # Metin kutuları için genel ayarlar
        margin = PPTXInches(0.5)
        title_height = PPTXInches(0.8) # Slayt başlığı için ayrılan yükseklik
        content_area_top = title_height + margin
        content_area_height = slide_height - content_area_top - margin
        content_area_width = slide_width - 2 * margin

        box_width = (content_area_width - margin) / 2 # İki sütun için
        box_height = (content_area_height - margin) / 2 # İki satır için

        for i, key_name in enumerate(swot_keys_ordered):
            config = swot_config[key_name]
            title = key_name
            color = config['color']
            items = swot_data_dict.get(config['key'], [])

            row = i // 2
            col = i % 2
            
            left = margin + col * (box_width + margin)
            top = content_area_top + row * (box_height + margin)

            try:
                txBox = slide.shapes.add_textbox(left, top, box_width, box_height)
                tf = txBox.text_frame
                tf.clear()
                tf.word_wrap = True

                # Başlık paragrafı
                p_title = tf.add_paragraph()
                p_title.text = title
                p_title.font.bold = True
                p_title.font.color.rgb = color
                p_title.font.size = PPTXPt(16) # Başlık boyutu büyütüldü
                p_title.space_after = PPTXPt(8)

                # İçerik maddeleri
                if items:
                    for item_text in items:
                        p_item = tf.add_paragraph()
                        p_item.text = f"• {item_text}" # Madde işareti eklendi
                        p_item.font.size = PPTXPt(11) # İçerik boyutu
                        p_item.level = 0 # Girinti için seviye ayarı (0 ana seviye)
                        p_item.space_after = PPTXPt(3)
                else:
                     p_item = tf.add_paragraph()
                     p_item.text = "Veri bulunamadı."
                     p_item.font.size = PPTXPt(11)
                     p_item.font.italic = True

            except Exception as e:
                logging.error(f"PPTX SWOT metin kutusu eklenemedi ({title}): {e}")
                txBox_err = slide.shapes.add_textbox(left, top, box_width, box_height)
                tf_err = txBox_err.text_frame; p_err = tf_err.add_paragraph()
                p_err.text = f"{title} yüklenemedi."; p_err.font.size = PPTXPt(10); p_err.font.color.rgb = RGBColor(255,0,0)


    def _add_pptx_analiz_ozeti_slide(self, prs):
        """PowerPoint sunumuna analiz özetini içeren slayt ekler."""
        slide = self._add_pptx_slide_with_title(prs, "Analiz Özeti")
        if not slide: return

        ozet_sections_pptx = {
            'Temel Değerlendirme': self.analiz_ozeti.get('temel_ozet', 'Değerlendirme bulunamadı.'),
            'Yatırım Değerlendirmesi': self.analiz_ozeti.get('yatirim_ozet', 'Değerlendirme bulunamadı.'),
            'Uygunluk Puanı Özeti': self.analiz_ozeti.get('uygunluk_ozet', 'Uygunluk puanı özeti bulunamadı.'),
            'Öneriler ve Tavsiyeler': self.analiz_ozeti.get('tavsiyeler', 'Tavsiye bulunamadı.')
        }

        # Slaytın ana içerik alanı için placeholder bulmaya çalışalım
        # Genellikle title'dan sonraki placeholder olur (layout_index=1 için)
        content_placeholder = None
        if len(slide.placeholders) > 1: # Başlık placeholder'ı dışında bir placeholder varsa
            for ph in slide.placeholders:
                if ph != slide.shapes.title and hasattr(ph, 'text_frame'): # Başlık değilse ve metin çerçevesi varsa
                    content_placeholder = ph
                    break
        
        if content_placeholder:
            tf = content_placeholder.text_frame
            tf.clear() # Mevcut içeriği temizle
            tf.word_wrap = True

            for title, text in ozet_sections_pptx.items():
                p_title = tf.add_paragraph()
                p_title.text = title
                p_title.font.bold = True
                p_title.font.size = PPTXPt(18) # Özet başlık boyutu
                p_title.space_after = PPTXPt(5)

                # Metni satır sonlarına göre böl ve her birini ayrı paragraf olarak ekle
                text_lines = str(text).splitlines()
                for line in text_lines:
                    if line.strip(): # Boş satırları atla
                        p_content = tf.add_paragraph()
                        p_content.text = line
                        p_content.font.size = PPTXPt(12) # Özet içerik boyutu
                        p_content.level = 1 # Biraz girintili
                        p_content.space_after = PPTXPt(3)
                
                # Bölümler arası boşluk için ek bir boş paragraf veya space_after ayarı
                if tf.paragraphs: # Eğer paragraf eklendiyse sonuncusuna boşluk ekle
                    tf.paragraphs[-1].space_after = PPTXPt(12)

        else: # Placeholder bulunamazsa manuel metin kutusu
            logging.warning("PPTX Analiz Özeti için içerik placeholder bulunamadı. Manuel metin kutusu kullanılıyor.")
            y_offset = PPTXInches(1.5)
            left = PPTXInches(0.5); width = PPTXInches(9)

            for title, text in ozet_sections_pptx.items():
                try:
                    txBox_title = slide.shapes.add_textbox(left, y_offset, width, PPTXInches(0.5))
                    tf_title = txBox_title.text_frame; p_title = tf_title.add_paragraph()
                    p_title.text = title; p_title.font.bold = True; p_title.font.size = PPTXPt(18)
                    y_offset += PPTXInches(0.5)

                    content_lines = str(text).splitlines()
                    estimated_height = PPTXInches(0.3 * len(content_lines) if content_lines else 0.3) # Yüksekliği tahmin et
                    txBox_content = slide.shapes.add_textbox(left, y_offset, width, estimated_height if estimated_height > PPTXInches(0.3) else PPTXInches(0.3))
                    tf_content = txBox_content.text_frame; tf_content.word_wrap = True
                    for line in content_lines:
                        if line.strip():
                            p_content = tf_content.add_paragraph()
                            p_content.text = line; p_content.font.size = PPTXPt(12)
                    y_offset += (estimated_height if estimated_height > PPTXInches(0.3) else PPTXInches(0.3)) + PPTXInches(0.3)
                except Exception as e:
                     logging.error(f"PPTX özet bölümü eklenemedi ({title}): {e}")
                     y_offset += PPTXInches(1.0)


    def _add_pptx_photos_slide(self, prs):
        """PowerPoint sunumuna yüklenen fotoğrafları içeren slayt ekler."""
        logging.info("_add_pptx_photos_slide metodu başladı.")
        # Fotoğraflar için "İki İçerik" veya "Boş" gibi bir layout daha uygun olabilir.
        # Şimdilik "Başlık ve İçerik" (layout 1) veya "Boş" (layout 5/6) kullanmaya devam edelim.
        slide = self._add_pptx_slide_with_title(prs, "Arsa Fotoğrafları", layout_index=5) # Boş layout deneyelim
        if not slide:
            logging.error("Fotoğraf slaytı oluşturulamadı.")
            return

        images = self._get_uploaded_images()
        logging.info(f"PPTX'e eklenecek resim sayısı: {len(images)}")

        if not images:
            txBox = slide.shapes.add_textbox(PPTXInches(1), PPTXInches(1.5), PPTXInches(8), PPTXInches(1))
            tf = txBox.text_frame; p = tf.add_paragraph()
            p.text = "Bu arsa için fotoğraf bulunmamaktadır."; p.font.size = PPTXPt(12); p.alignment = 2 # CENTER
            logging.info("Eklenecek resim bulunamadı, uyarı mesajı eklendi.")
            return

        # Slayt boyutları
        slide_w = prs.slide_width
        slide_h = prs.slide_height

        # Resimleri slayta yerleştirme (basit ızgara düzeni)
        # En fazla 4 resim gösterelim (2x2 grid)
        num_images_to_show = min(len(images), 4)
        cols = 2 if num_images_to_show > 1 else 1
        rows = (num_images_to_show + cols - 1) // cols

        margin = PPTXInches(0.5)
        # Başlık için manuel bir alan bırakalım (layout boş olsa bile)
        title_space_top = PPTXInches(1.2) if slide.shapes.title is None else PPTXInches(0) # Eğer başlık manuel eklendiyse

        img_area_width = slide_w - 2 * margin
        img_area_height = slide_h - title_space_top - margin * 2 # Alt ve üst marj

        img_width = (img_area_width - (cols - 1) * margin) / cols
        img_height = (img_area_height - (rows - 1) * margin) / rows


        for i in range(num_images_to_show):
            img_path = images[i]
            row_idx = i // cols
            col_idx = i % cols

            left = margin + col_idx * (img_width + margin)
            top = title_space_top + margin + row_idx * (img_height + margin)
            
            # Resim boyutunu orantılı olarak ayarla
            try:
                with Image.open(img_path) as pil_img:
                    pil_w, pil_h = pil_img.size
                    aspect_ratio = pil_w / pil_h if pil_h > 0 else 1

                current_img_width = img_width
                current_img_height = img_width / aspect_ratio

                if current_img_height > img_height:
                    current_img_height = img_height
                    current_img_width = img_height * aspect_ratio
                
                # Ortalamak için pozisyonu ayarla (hücre içinde)
                actual_left = left + (img_width - current_img_width) / 2
                actual_top = top + (img_height - current_img_height) / 2


                slide.shapes.add_picture(img_path, actual_left, actual_top, width=current_img_width, height=current_img_height)
                logging.info(f"PPTX resim eklendi: {img_path}")
            except FileNotFoundError:
                logging.error(f"PPTX resim dosyası bulunamadı: {img_path}")
            except Exception as e:
                logging.error(f"PPTX resim eklenemedi: {img_path} - {e}")
                txBox = slide.shapes.add_textbox(left, top, img_width, PPTXInches(0.5)) # Hata mesajı için küçük kutu
                tf = txBox.text_frame; p = tf.add_paragraph()
                p.text = "Resim Hatalı"; p.font.size = PPTXPt(10); p.alignment = 2; p.font.color.rgb = RGBColor(255,0,0)

        logging.info("_add_pptx_photos_slide metodu tamamlandı.")


    def _add_pptx_qr_code_slide(self, prs, qr_path):
        """PowerPoint sunumuna QR kod slaytı ekler."""
        slide = self._add_pptx_slide_with_title(prs, "Analiz Detayları QR Kod", layout_index=5) # Boş layout
        if not slide: return

        # QR kod resmini slaytın ortasına ekle
        if os.path.exists(qr_path):
            try:
                slide_width = prs.slide_width
                slide_height = prs.slide_height

                qr_dim = PPTXInches(3.5) # QR kod boyutu (kare)

                left = (slide_width - qr_dim) / 2
                # Başlık için manuel bir alan bırakalım
                title_space_top = PPTXInches(1.2) if slide.shapes.title is None else PPTXInches(0)
                top_qr = title_space_top + (slide_height - title_space_top - qr_dim - PPTXInches(0.8)) / 2 # QR'ı dikeyde ortala (açıklama için yer bırak)


                slide.shapes.add_picture(qr_path, left, top_qr, width=qr_dim, height=qr_dim)
                logging.info(f"PPTX QR kod resmi eklendi: {qr_path}")

                # Açıklama ekle (QR kodun altına)
                caption_text = "Bu QR kodu tarayarak analiz detaylarına ulaşabilirsiniz."
                cap_left = left; cap_top = top_qr + qr_dim + PPTXInches(0.2)
                cap_width = qr_dim; cap_height = PPTXInches(0.6)

                txBox_caption = slide.shapes.add_textbox(cap_left, cap_top, cap_width, cap_height)
                tf_caption = txBox_caption.text_frame; tf_caption.clear()
                p_caption = tf_caption.add_paragraph()
                p_caption.text = caption_text
                p_caption.font.size = PPTXPt(12) # Açıklama boyutu
                p_caption.alignment = 2 # CENTER
                p_caption.font.italic = True
                p_caption.font.color.rgb = RGBColor(100, 100, 100) # Gri tonu

            except Exception as e:
                logging.error(f"PPTX QR kod resmi eklenemedi: {e}")
                txBox = slide.shapes.add_textbox(PPTXInches(3.5), PPTXInches(3), PPTXInches(3), PPTXInches(1))
                tf = txBox.text_frame; p = tf.add_paragraph()
                p.text = "QR Kod yüklenemedi."; p.font.size = PPTXPt(10); p.alignment = 2; p.font.color.rgb = RGBColor(255,0,0)
        else:
             logging.warning("QR kod dosyası bulunamadı, PPTX'e eklenemedi.")
             txBox = slide.shapes.add_textbox(PPTXInches(3.5), PPTXInches(3), PPTXInches(3), PPTXInches(1))
             tf = txBox.text_frame; p = tf.add_paragraph()
             p.text = "QR Kod Yüklenemedi (Dosya Bulunamadı)."; p.font.size = PPTXPt(10); p.alignment = 2; p.font.color.rgb = RGBColor(255,0,0)

