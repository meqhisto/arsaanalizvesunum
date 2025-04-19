"""
Arsa Yatırım Danışmanlığı - Sunum Oluşturucu Modülü
Bu modül, arsa analiz sonuçlarından Word ve PowerPoint sunumları oluşturmak için kullanılır.
"""

import os
from docx import Document
from docx.shared import Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN

class SunumOlusturucu:
    def __init__(self, templates_dir):
        self.templates_dir = templates_dir
        os.makedirs(templates_dir, exist_ok=True)
        
        self.word_template = os.path.join(templates_dir, "word_template.docx")
        self.ppt_template = os.path.join(templates_dir, "ppt_template.pptx")
        
        self._create_templates_if_not_exist()
    
    def _create_templates_if_not_exist(self):
        if not os.path.exists(self.word_template):
            self._create_word_template()
        if not os.path.exists(self.ppt_template):
            self._create_ppt_template()
    
    def _create_word_template(self):
        doc = Document()
        
        for section in doc.sections:
            section.page_height = Cm(29.7)
            section.page_width = Cm(21)
            section.left_margin = Cm(2.5)
            section.right_margin = Cm(2.5)
            section.top_margin = Cm(2.5)
            section.bottom_margin = Cm(2.5)
        
        title = doc.add_heading("ARSA YATIRIM ANALİZ RAPORU", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        doc.add_heading("1. ARSA BİLGİLERİ", level=1)
        doc.add_paragraph("Konum: {konum}")
        doc.add_paragraph("Metrekare: {metrekare} m²")
        doc.add_paragraph("İmar Durumu: {imar_durumu}")
        doc.add_paragraph("Fiyat: {fiyat} TL")
        doc.add_paragraph("Metrekare Fiyatı: {metrekare_fiyat} TL/m²")
        doc.add_paragraph("Bölge Ortalama Fiyatı: {bolge_fiyat} TL/m²")
        
        doc.add_heading("2. ANALİZ SONUÇLARI", level=1)
        doc.add_paragraph("Bölge Karşılaştırması: {bolge_karsilastirma}%")
        doc.add_paragraph("Potansiyel Getiri: %{potansiyel_getiri}")
        doc.add_paragraph("Tavsiye Edilen Yatırım Süresi: {yatirim_suresi} yıl")
        
        doc.add_heading("3. YATIRIM DEĞERLENDİRMESİ", level=1)
        doc.add_paragraph("{temel_ozet}")
        doc.add_paragraph("{yatirim_ozet}")
        
        doc.add_heading("4. TAVSİYELER", level=1)
        doc.add_paragraph("{tavsiyeler}")
        
        doc.save(self.word_template)
    
    def _create_ppt_template(self):
        prs = Presentation()
        
        # Kapak slaytı
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "ARSA YATIRIM ANALİZ RAPORU"
        slide.placeholders[1].text = "{konum}"
        
        # Diğer slaytlar
        layouts = [
            ("ARSA BİLGİLERİ", "• Konum: {konum}\n• Metrekare: {metrekare} m²\n• İmar Durumu: {imar_durumu}\n• Fiyat: {fiyat} TL\n• Metrekare Fiyatı: {metrekare_fiyat} TL/m²\n• Bölge Ortalama Fiyatı: {bolge_fiyat} TL/m²"),
            ("ANALİZ SONUÇLARI", "• Bölge Karşılaştırması: {bolge_karsilastirma}%\n• Potansiyel Getiri: %{potansiyel_getiri}\n• Tavsiye Edilen Yatırım Süresi: {yatirim_suresi} yıl"),
            ("YATIRIM DEĞERLENDİRMESİ", "{temel_ozet}\n\n{yatirim_ozet}"),
            ("TAVSİYELER", "{tavsiyeler}")
        ]
        
        for title, content in layouts:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = title
            slide.placeholders[1].text = content
        
        prs.save(self.ppt_template)
    
    def _format_number(self, number, decimal_places=2):
        formatted = f"{number:,.{decimal_places}f}".replace(",", "X").replace(".", ",").replace("X", ".")
        return formatted
    
    def olustur_word(self, arsa_data, output_path):
        doc = Document(self.word_template)
        
        konum = f"{arsa_data['konum']['mahalle']}, {arsa_data['konum']['ilce']}, {arsa_data['konum']['il']}"
        replacements = {
            "konum": konum,
            "metrekare": self._format_number(arsa_data['metrekare']),
            "imar_durumu": arsa_data['imar_durumu'].capitalize(),
            "fiyat": self._format_number(arsa_data['fiyat']),
            "metrekare_fiyat": self._format_number(arsa_data['metrekare_fiyat']),
            "bolge_fiyat": self._format_number(arsa_data['bolge_fiyat']),
            "bolge_karsilastirma": self._format_number(arsa_data['bolge_karsilastirma'], 1),
            "potansiyel_getiri": self._format_number(arsa_data['potansiyel_getiri'], 1),
            "yatirim_suresi": str(arsa_data['yatirim_suresi']),
            "temel_ozet": arsa_data.get('ozet', {}).get('temel_ozet', ''),
            "yatirim_ozet": arsa_data.get('ozet', {}).get('yatirim_ozet', ''),
            "tavsiyeler": arsa_data.get('ozet', {}).get('tavsiyeler', '')
        }
        
        for paragraph in doc.paragraphs:
            for key, value in replacements.items():
                paragraph.text = paragraph.text.replace(f"{{{key}}}", value)
        
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        doc.save(output_path)
        return output_path
    
    def olustur_powerpoint(self, arsa_data, output_path):
        prs = Presentation(self.ppt_template)
        
        konum = f"{arsa_data['konum']['mahalle']}, {arsa_data['konum']['ilce']}, {arsa_data['konum']['il']}"
        replacements = {
            "konum": konum,
            "metrekare": self._format_number(arsa_data['metrekare']),
            "imar_durumu": arsa_data['imar_durumu'].capitalize(),
            "fiyat": self._format_number(arsa_data['fiyat']),
            "metrekare_fiyat": self._format_number(arsa_data['metrekare_fiyat']),
            "bolge_fiyat": self._format_number(arsa_data['bolge_fiyat']),
            "bolge_karsilastirma": self._format_number(arsa_data['bolge_karsilastirma'], 1),
            "potansiyel_getiri": self._format_number(arsa_data['potansiyel_getiri'], 1),
            "yatirim_suresi": str(arsa_data['yatirim_suresi']),
            "temel_ozet": arsa_data.get('ozet', {}).get('temel_ozet', ''),
            "yatirim_ozet": arsa_data.get('ozet', {}).get('yatirim_ozet', ''),
            "tavsiyeler": arsa_data.get('ozet', {}).get('tavsiyeler', '')
        }
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    for key, value in replacements.items():
                        shape.text = shape.text.replace(f"{{{key}}}", value)
        
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        prs.save(output_path)
        return output_path